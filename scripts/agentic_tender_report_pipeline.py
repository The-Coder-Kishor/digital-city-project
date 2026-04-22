from __future__ import annotations

import argparse
import html as pyhtml
import json
import os
import re
import shutil
import subprocess
import time
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

import requests
from lxml import etree, html as lxml_html

from telangana_tenders import TelanganaTenderClient, parse_html_to_json, recursive_unzip
from telangana_tenders.client import validate_zip

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".tsv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".log",
    ".yaml",
    ".yml",
    ".rtf",
}

SKIP_BINARY_EXTENSIONS = {
    ".zip",
    ".7z",
    ".rar",
    ".exe",
    ".dll",
    ".jpg",
    ".jpeg",
    ".png",
    ".gif",
    ".bmp",
    ".webp",
    ".mp4",
    ".mp3",
}


@dataclass(slots=True)
class ExtractionIssue:
    file: str
    extension: str
    reason: str


@dataclass(slots=True)
class SearchResult:
    title: str
    href: str
    snippet: str
    source_query: str


class OllamaChatClient:
    def __init__(
        self,
        model: str,
        base_url: str = "http://localhost:11434",
        timeout_seconds: float = 300.0,
    ) -> None:
        self.base_url = base_url.rstrip("/")
        self.model = model
        self.timeout_seconds = timeout_seconds
        self._session = requests.Session()

    def chat(
        self,
        messages: list[dict[str, str]],
        *,
        temperature: float = 0.0,
        max_tokens: int = 2000,
        num_ctx: int | None = None,
    ) -> str:
        payload: dict[str, Any] = {
            "model": self.model,
            "messages": messages,
            "stream": False,
            "options": {
                "temperature": temperature,
                "num_predict": max_tokens,
                "num_ctx": int(num_ctx) if num_ctx is not None else 16384,
            },
        }
        resp = self._session.post(
            f"{self.base_url}/api/chat",
            json=payload,
            headers={"Content-Type": "application/json"},
            timeout=self.timeout_seconds,
        )
        resp.raise_for_status()
        content = ((resp.json() or {}).get("message") or {}).get("content")
        if not isinstance(content, str) or not content.strip():
            raise ValueError("Ollama returned empty content")
        return content


class DocumentReader:
    def __init__(self, *, max_files: int = 200, max_chars_per_file: int = 12000) -> None:
        self.max_files = max_files
        self.max_chars_per_file = max_chars_per_file

    def read_documents(self, source_dirs: list[Path]) -> tuple[list[dict[str, str]], list[ExtractionIssue]]:
        docs: list[dict[str, str]] = []
        issues: list[ExtractionIssue] = []

        files: list[Path] = []
        for source_dir in source_dirs:
            if source_dir.exists():
                files.extend([p for p in source_dir.rglob("*") if p.is_file()])

        files.sort(key=lambda p: str(p).lower())
        files = files[: self.max_files]

        for path in files:
            text, handled, reason = self._extract_text(path)
            normalized = self._normalize_text(text)

            if reason:
                issues.append(
                    ExtractionIssue(
                        file=str(path).replace("\\", "/"),
                        extension=path.suffix.lower(),
                        reason=reason,
                    )
                )

            if not handled or not normalized:
                continue

            docs.append(
                {
                    "path": str(path).replace("\\", "/"),
                    "text": normalized[: self.max_chars_per_file],
                }
            )

        return docs, issues

    def _extract_text(self, path: Path) -> tuple[str, bool, str | None]:
        suffix = path.suffix.lower()

        if suffix in SKIP_BINARY_EXTENSIONS:
            return "", True, None

        if suffix in TEXT_EXTENSIONS:
            try:
                raw = path.read_text(encoding="utf-8", errors="ignore")
                if suffix == ".rtf":
                    return self._extract_rtf_text(raw), True, None
                return raw, True, None
            except OSError:
                return "", True, f"Failed to read text file: {path.name}"

        if suffix in {".docx", ".docm"}:
            text = self._extract_docx_text(path)
            return (text, True, None) if text else ("", True, f"Unable to extract {suffix} text")

        if suffix == ".pdf":
            text = self._extract_pdf_text(path)
            return (text, True, None) if text else ("", True, "Unable to extract PDF text")

        if suffix == ".doc":
            text = self._extract_doc_text(path)
            return (
                (text, True, "Legacy DOC extracted heuristically; verify confidence")
                if text
                else ("", True, "Unable to extract DOC text")
            )

        if suffix in {".xlsx", ".xlsm"}:
            text = self._extract_xlsx_text(path)
            return (text, True, None) if text else ("", True, f"Unable to extract {suffix} text")

        if suffix == ".xls":
            text = self._extract_xls_text(path)
            return (text, True, None) if text else ("", True, "Unable to extract XLS text")

        if suffix == ".pptx":
            text = self._extract_pptx_text(path)
            return (text, True, None) if text else ("", True, "Unable to extract PPTX text")

        if suffix == ".ppt":
            text = self._extract_printable_runs(path)
            return (
                (text, True, "Legacy PPT extracted heuristically; verify confidence")
                if text
                else ("", True, "Unable to extract PPT text")
            )

        if suffix in {".odt", ".ods"}:
            text = self._extract_open_document_text(path)
            return (text, True, None) if text else ("", True, f"Unable to extract {suffix} text")

        return "", False, f"Unhandled file extension: {suffix or '<none>'}"

    def _extract_rtf_text(self, text: str) -> str:
        text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
        text = re.sub(r"\\[a-zA-Z]+-?\d*\s?", " ", text)
        return text.replace("{", " ").replace("}", " ")

    def _extract_docx_text(self, path: Path) -> str:
        try:
            import docx  # type: ignore

            document = docx.Document(str(path))
            paragraphs = [p.text for p in document.paragraphs if p.text]
            if paragraphs:
                return "\n".join(paragraphs)
        except Exception:
            pass

        try:
            with zipfile.ZipFile(path) as archive:
                xml_data = archive.read("word/document.xml")
            root = etree.fromstring(xml_data)
            return " ".join(root.itertext())
        except Exception:
            return ""

    def _extract_pdf_text(self, path: Path) -> str:
        try:
            from pypdf import PdfReader  # type: ignore

            reader = PdfReader(str(path))
            parts = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    parts.append(page_text)
            if parts:
                return "\n".join(parts)
        except Exception:
            pass

        return self._extract_printable_runs(path)

    def _extract_doc_text(self, path: Path) -> str:
        antiword = shutil.which("antiword")
        if antiword:
            try:
                completed = subprocess.run(
                    [antiword, str(path)],
                    check=False,
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="ignore",
                )
                if completed.returncode == 0 and completed.stdout.strip():
                    return completed.stdout
            except OSError:
                pass

        return self._extract_printable_runs(path)

    def _extract_xlsx_text(self, path: Path) -> str:
        try:
            from openpyxl import load_workbook  # type: ignore

            workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in workbook.worksheets[:12]:
                parts.append(f"[Sheet] {sheet.title}")
                for row in sheet.iter_rows(min_row=1, max_row=300):
                    values = [str(cell.value).strip() for cell in row if cell.value is not None]
                    if values:
                        parts.append(" | ".join(values))
            return "\n".join(parts)
        except Exception:
            return ""

    def _extract_xls_text(self, path: Path) -> str:
        try:
            import xlrd  # type: ignore

            book = xlrd.open_workbook(str(path), on_demand=True)
            parts: list[str] = []
            for sheet in book.sheets()[:10]:
                parts.append(f"[Sheet] {sheet.name}")
                for row_idx in range(min(sheet.nrows, 300)):
                    row_values = [str(v).strip() for v in sheet.row_values(row_idx) if str(v).strip()]
                    if row_values:
                        parts.append(" | ".join(row_values))
            return "\n".join(parts)
        except Exception:
            return ""

    def _extract_pptx_text(self, path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore

            presentation = Presentation(str(path))
            parts: list[str] = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                parts.append(f"[Slide {slide_index}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        value = str(shape.text).strip()
                        if value:
                            parts.append(value)
            return "\n".join(parts)
        except Exception:
            return ""

    def _extract_open_document_text(self, path: Path) -> str:
        try:
            with zipfile.ZipFile(path) as archive:
                xml_data = archive.read("content.xml")
            root = etree.fromstring(xml_data)
            return " ".join(root.itertext())
        except Exception:
            return ""

    def _extract_printable_runs(self, path: Path) -> str:
        try:
            data = path.read_bytes()
        except OSError:
            return ""
        candidate = data.decode("latin-1", errors="ignore")
        lines = re.findall(r"[A-Za-z0-9][^\r\n]{20,}", candidate)
        return "\n".join(lines[:500])

    def _normalize_text(self, text: str) -> str:
        if not text:
            return ""
        if "<html" in text.lower() or "<table" in text.lower():
            try:
                tree = lxml_html.fromstring(text)
                text = tree.text_content()
            except Exception:
                pass
        return "\n".join(line.strip() for line in text.splitlines() if line.strip())


class LocalNewsCrawler:
    def __init__(self, *, timeout_seconds: float = 15.0, max_workers: int = 6) -> None:
        self.timeout_seconds = timeout_seconds
        self.max_workers = max_workers
        self.session = requests.Session()
        self.session.headers.update(
            {
                "User-Agent": (
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                    "(KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36"
                )
            }
        )

    def search(self, queries: list[str], *, max_results_per_query: int = 10) -> list[SearchResult]:
        DDGS = None
        try:
            from ddgs import DDGS as DDGSClient  # type: ignore

            DDGS = DDGSClient
        except Exception:
            try:
                # Backward compatibility for older environments.
                from duckduckgo_search import DDGS as DDGSClient  # type: ignore

                DDGS = DDGSClient
            except Exception as exc:
                raise RuntimeError(
                    "ddgs is required for local web crawling. Install dependencies with `uv sync`."
                ) from exc

        assert DDGS is not None

        def _safe_text_search(ddgs_client: Any, query: str) -> list[dict[str, Any]]:
            try:
                rows = ddgs_client.text(query, max_results=max_results_per_query) or []
                return [row for row in rows if isinstance(row, dict)]
            except Exception:
                # ddgs raises exceptions (e.g. "No results found") for empty results.
                # Treat those as a normal no-hit case and continue.
                return []

        seen: set[str] = set()
        results: list[SearchResult] = []
        with DDGS() as ddgs:
            for query in queries:
                rows = _safe_text_search(ddgs, query)
                for row in rows:
                    href = str(row.get("href") or "").strip()
                    if not href or href in seen:
                        continue
                    seen.add(href)
                    results.append(
                        SearchResult(
                            title=str(row.get("title") or "").strip(),
                            href=href,
                            snippet=str(row.get("body") or "").strip(),
                            source_query=query,
                        )
                    )

        # Fallback broadening when highly specific queries return no usable results.
        if not results:
            fallback_queries = []
            for query in queries:
                q = query.replace('"', "").strip()
                if q:
                    fallback_queries.append(q)
            with DDGS() as ddgs:
                for query in fallback_queries[:6]:
                    rows = _safe_text_search(ddgs, query)
                    for row in rows:
                        href = str(row.get("href") or "").strip()
                        if not href or href in seen:
                            continue
                        seen.add(href)
                        results.append(
                            SearchResult(
                                title=str(row.get("title") or "").strip(),
                                href=href,
                                snippet=str(row.get("body") or "").strip(),
                                source_query=query,
                            )
                        )
        return results

    def crawl(self, results: list[SearchResult], *, max_pages: int = 10, max_chars: int = 14000) -> list[dict[str, str]]:
        selected = results[:max_pages]
        pages: list[dict[str, str]] = []

        def fetch_one(item: SearchResult) -> dict[str, str]:
            try:
                resp = self.session.get(item.href, timeout=self.timeout_seconds)
                resp.raise_for_status()
                content_type = (resp.headers.get("content-type") or "").lower()
                if "text/html" not in content_type and "xml" not in content_type:
                    return {
                        "url": item.href,
                        "title": item.title,
                        "query": item.source_query,
                        "snippet": item.snippet,
                        "content": "",
                        "error": f"Skipped non-HTML content type: {content_type}",
                    }

                tree = lxml_html.fromstring(resp.text)
                for bad in tree.xpath("//script|//style|//noscript"):
                    bad.getparent().remove(bad)
                text = " ".join(tree.text_content().split())
                return {
                    "url": item.href,
                    "title": item.title,
                    "query": item.source_query,
                    "snippet": item.snippet,
                    "content": text[:max_chars],
                    "error": "",
                }
            except Exception as exc:
                return {
                    "url": item.href,
                    "title": item.title,
                    "query": item.source_query,
                    "snippet": item.snippet,
                    "content": "",
                    "error": f"crawl failed: {exc}",
                }

        with ThreadPoolExecutor(max_workers=max(1, self.max_workers)) as executor:
            futures = [executor.submit(fetch_one, item) for item in selected]
            for future in as_completed(futures):
                pages.append(future.result())

        return pages


def extract_json_object(text: str) -> dict[str, Any]:
    candidate = text.strip()
    if candidate.startswith("```"):
        candidate = re.sub(r"^```[a-zA-Z0-9_\-]*", "", candidate).strip()
        if candidate.endswith("```"):
            candidate = candidate[:-3].strip()
    try:
        return json.loads(candidate)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", text, flags=re.DOTALL)
        if not match:
            raise
        return json.loads(match.group(0))


def write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=True), encoding="utf-8")


def write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def _parse_details_kv(details: list[Any]) -> dict[str, str]:
    out: dict[str, str] = {}
    for line in details:
        if not isinstance(line, str) or ":" not in line:
            continue
        key, _, value = line.partition(":")
        k = key.strip()
        v = value.strip()
        if k and v:
            out[k.lower()] = v
    return out


def _extract_procurement_refs(text: str) -> list[str]:
    if not text:
        return []
    patterns = [
        r"\b\d{1,4}/[A-Z0-9()\-/]+/[A-Z][A-Z0-9()\-/]+/\d{4}-\d{2,4}\b",
        r"\bNIT[/\s-]*No\.?\s*[:\-]?\s*[^\s,;]+",
        r"\bGEM/\d{4,}/[A-Z0-9/\-]+\b",
    ]
    found: list[str] = []
    for pat in patterns:
        for m in re.finditer(pat, text, flags=re.IGNORECASE):
            s = m.group(0).strip()
            if s and s not in found:
                found.append(s)
    return found[:12]


def _short_work_title(title: str, *, max_words: int = 14) -> str:
    t = re.sub(r"\s+", " ", (title or "").strip())
    if not t:
        return ""
    words = t.split()
    if len(words) <= max_words:
        return t
    return " ".join(words[:max_words])


def _geo_tokens_from_title(title: str) -> list[str]:
    t = title or ""
    hits = re.findall(
        r"\b(?:Circle\s+\d+[A-Z]?|KPHB|Kukatpally|GHMC|Secunderabad|Hyderabad|"
        r"Telangana|Warangal|Karimnagar|Nizamabad|Medchal|Rangareddy)\b",
        t,
        flags=re.IGNORECASE,
    )
    seen: set[str] = set()
    out: list[str] = []
    for h in hits:
        v = h.strip()
        lk = v.lower()
        if lk not in seen:
            seen.add(lk)
            out.append(v)
    return out


def merge_project_intel(doc_analysis: dict[str, Any]) -> dict[str, Any]:
    """
    Normalise project/work context for web search and deep research.
    Prefer LLM `project_intel`; fill gaps from `details` and timeline text.
    """
    raw = doc_analysis.get("project_intel")
    intel: dict[str, Any] = dict(raw) if isinstance(raw, dict) else {}

    details = doc_analysis.get("details") or []
    kv = _parse_details_kv([d for d in details if isinstance(d, str)])

    title = str(intel.get("title_or_work") or kv.get("name of work") or "").strip()
    department = str(intel.get("department") or kv.get("department name") or "").strip()
    authority = str(
        intel.get("authority_or_circle")
        or kv.get("circle/division")
        or kv.get("circle / division")
        or ""
    ).strip()
    work_type = str(intel.get("work_type") or kv.get("type of work") or "").strip()
    category = str(intel.get("category") or kv.get("tender category") or "").strip()

    timeline = doc_analysis.get("timeline") or []
    timeline_blob = " ".join(
        f"{item.get('event','')} {item.get('date','')} {item.get('evidence','')}"
        for item in timeline
        if isinstance(item, dict)
    )
    details_blob = " ".join(str(d) for d in details if isinstance(d, str))
    nit_from_intel = str(intel.get("nit_or_reference") or "").strip()
    refs = _extract_procurement_refs(" ".join(filter(None, [nit_from_intel, timeline_blob, details_blob])))

    locations = intel.get("locations")
    if not isinstance(locations, list):
        locations = []
    locations = [str(x).strip() for x in locations if isinstance(x, str) and str(x).strip()]
    for g in _geo_tokens_from_title(title):
        if g not in locations:
            locations.append(g)

    year_hint = str(intel.get("year_hint") or "").strip()
    if not year_hint:
        ym = re.search(r"\b(20\d{2})\b", timeline_blob + " " + details_blob)
        if ym:
            year_hint = ym.group(1)

    intel.setdefault("title_or_work", title)
    intel.setdefault("department", department)
    intel.setdefault("authority_or_circle", authority)
    intel.setdefault("work_type", work_type)
    intel.setdefault("category", category)
    intel.setdefault("nit_or_reference", nit_from_intel or (refs[0] if refs else ""))
    intel.setdefault("procurement_reference_codes", refs)
    intel.setdefault("locations", locations)
    intel.setdefault("year_hint", year_hint)

    sq = intel.get("search_queries")
    if not isinstance(sq, list):
        sq = []
    intel["search_queries"] = [str(x).strip() for x in sq if isinstance(x, str) and str(x).strip()]
    return intel


def build_web_queries(
    tender_id: str,
    doc_analysis: dict[str, Any],
    *,
    project_intel: dict[str, Any],
    max_queries: int = 42,
) -> list[str]:
    """
    Prefer project/work-specific queries (title, NIT, department, geography).
    Bare 6-digit portal IDs are poor search keys; when used, always pair with disambiguators.
    """
    intel = project_intel
    title = str(intel.get("title_or_work") or "").strip()
    short_title = _short_work_title(title)
    dept = str(intel.get("department") or "").strip()
    authority = str(intel.get("authority_or_circle") or "").strip()
    work_type = str(intel.get("work_type") or "").strip()
    year = str(intel.get("year_hint") or "").strip()
    nit = str(intel.get("nit_or_reference") or "").strip()
    refs = intel.get("procurement_reference_codes") or []
    ref_list = [str(r).strip() for r in refs if isinstance(r, str) and str(r).strip()]
    locations = [str(x).strip() for x in (intel.get("locations") or []) if str(x).strip()]
    geo_phrase = " ".join(locations[:4]) if locations else ""

    queries: list[str] = []

    # Model-suggested queries first (should already be project-centric).
    for q in intel.get("search_queries") or []:
        if isinstance(q, str) and q.strip():
            queries.append(q.strip())

    if nit:
        queries.append(nit)
        queries.append(f'"{nit}" tender')
        queries.append(f'"{nit}" GHMC')
    for ref in ref_list[:6]:
        if ref != nit:
            queries.append(ref)

    if short_title:
        queries.append(f'"{short_title}"')
        queries.append(f"{short_title} tender Telangana")
        if dept:
            queries.append(f"{short_title} {dept}")
        if year:
            queries.append(f"{short_title} {year}")

    if dept and work_type:
        queries.append(f"{dept} {work_type} {year}".strip())
    if dept and geo_phrase:
        queries.append(f"{dept} {geo_phrase} {work_type}".strip())

    if title and "GHMC" in title.upper():
        queries.append(f"site:tender.telangana.gov.in GHMC {_short_work_title(title, max_words=8)}")

    # Portal / news angles (still project words, not ID-only).
    if short_title:
        queries.append(f"{short_title} corrigendum")
        queries.append(f"{short_title} award")
        queries.append(f"{short_title} eprocurement Telangana")

    # Disambiguated portal ID queries (last resort, not first).
    disamb = " ".join(x for x in (dept, geo_phrase, year, "Telangana") if x).strip()
    if disamb:
        queries.append(f"{tender_id} {disamb} tender")
        queries.append(f"site:tender.telangana.gov.in {tender_id} {dept.split()[0] if dept else ''}".strip())
    queries.append(f'"{tender_id}" {short_title or dept or "Telangana"}')

    seen: set[str] = set()
    unique: list[str] = []
    for q in queries:
        qn = re.sub(r"\s+", " ", q).strip()
        if not qn or qn in seen:
            continue
        seen.add(qn)
        unique.append(qn)
        if len(unique) >= max_queries:
            break
    return unique


def collect_tender_data(tender_id: str, output_dir: Path, rate_limit_seconds: float) -> dict[str, Any]:
    tender_dir = output_dir / tender_id
    inspect_dir = tender_dir / "inspect"
    inspect_dir.mkdir(parents=True, exist_ok=True)

    client = TelanganaTenderClient(rate_limit_seconds=rate_limit_seconds)
    bootstrap_error = ""
    try:
        cookies = client.bootstrap_cookies()
    except Exception as exc:
        cookies = {}
        bootstrap_error = str(exc)

    portal_data: dict[str, Any] = {
        "cookies": cookies,
        "bootstrap_error": bootstrap_error,
        "endpoints": {},
    }

    endpoint_calls = [
        ("general", client.fetch_general_info),
        ("documents", client.fetch_documents_info),
        ("award", client.fetch_award_details),
        ("award_documents", client.fetch_award_documents),
    ]

    for label, func in endpoint_calls:
        try:
            page = func(tender_id)
            page_json = page.to_json()
            write_text(inspect_dir / f"{label}.html", page.html_text)
            write_json(inspect_dir / f"{label}.json", page_json)
            portal_data["endpoints"][label] = {
                "status": "ok",
                "url": page.url,
                "text_preview": page.text[:600],
            }
        except Exception as exc:
            portal_data["endpoints"][label] = {
                "status": "error",
                "error": str(exc),
            }

    zip_path = tender_dir / f"{tender_id}_documents.zip"
    extracted_dir = tender_dir / f"{tender_id}_extracted"
    download_status = {"zip_path": str(zip_path), "extracted_dir": str(extracted_dir), "downloaded": False}

    try:
        client.download_documents_zip(tender_id, zip_path)
        download_status["downloaded"] = True
        download_status["zip_valid"] = validate_zip(zip_path)
        if download_status["zip_valid"]:
            extracted_dir.mkdir(parents=True, exist_ok=True)
            archives = recursive_unzip(zip_path, extracted_dir)
            download_status["archives_extracted"] = len(archives)
        else:
            download_status["archives_extracted"] = 0
    except Exception as exc:
        download_status["error"] = str(exc)

    source_dirs: list[Path] = []
    if extracted_dir.exists() and any(extracted_dir.rglob("*")):
        source_dirs.append(extracted_dir)
    if inspect_dir.exists() and any(inspect_dir.rglob("*")):
        source_dirs.append(inspect_dir)

    return {
        "tender_dir": str(tender_dir),
        "inspect_dir": str(inspect_dir),
        "portal_data": portal_data,
        "download_status": download_status,
        "source_dirs": [str(p) for p in source_dirs],
    }


def analyze_documents(
    llm: OllamaChatClient,
    tender_id: str,
    docs: list[dict[str, str]],
    portal_data: dict[str, Any],
) -> dict[str, Any]:
    manifest = [{"path": d["path"], "chars": len(d["text"])} for d in docs]
    snippets = [{"path": d["path"], "snippet": d["text"][:2500]} for d in docs[:18]]

    messages = [
        {
            "role": "system",
            "content": (
                "You are a tender analyst. Extract key facts from tender portal data and documents. "
                "Return JSON only."
            ),
        },
        {
            "role": "user",
            "content": (
                f"Tender ID: {tender_id}\n"
                "Return JSON schema:\n"
                "{"
                "\"tender_id\":\"...\","
                "\"status\":\"awarded|not_awarded|unknown\","
                "\"awarded_by\":{\"value\":\"...\",\"evidence\":[\"...\"]}|null,"
                "\"awarded_to\":{\"value\":\"...\",\"evidence\":[\"...\"]}|null,"
                "\"details\":[\"...\"],"
                "\"timeline\":[{\"event\":\"...\",\"date\":\"...\",\"evidence\":\"...\"}],"
                "\"contract_value\":{\"amount\":\"...\",\"evidence\":[\"...\"]}|null,"
                "\"extra_points\":[{\"point\":\"...\",\"evidence\":\"...\"}],"
                "\"project_intel\":{"
                "\"title_or_work\":\"exact Name of Work / scope\","
                "\"department\":\"...\","
                "\"authority_or_circle\":\"...\","
                "\"nit_or_reference\":\"NIT / enquiry / schedule ref if any\","
                "\"locations\":[\"KPHB\",\"Kukatpally\", \"...\"],"
                "\"work_type\":\"...\","
                "\"category\":\"...\","
                "\"year_hint\":\"20xx\","
                "\"search_queries\":[\"12-22 web search strings for THIS project — use work title, department, "
                "NIT/ref, geography, work type; do NOT use only the 6-digit portal tender id\"]"
                "},"
                "\"confidence\":0.0,"
                "\"gaps\":[\"...\"],"
                "\"sources_used\":[\"...\"]"
                "}\n\n"
                f"Portal data summary:\n{json.dumps(portal_data, ensure_ascii=True)[:20000]}\n\n"
                f"Document manifest:\n{json.dumps(manifest, ensure_ascii=True)}\n\n"
                f"Document snippets:\n{json.dumps(snippets, ensure_ascii=True)}"
            ),
        },
    ]

    try:
        output = llm.chat(messages, temperature=0.0, max_tokens=3200, num_ctx=32768)
        payload = extract_json_object(output)
        if not isinstance(payload, dict):
            raise ValueError("document analysis returned non-object JSON")
    except Exception as exc:
        payload = {
            "tender_id": tender_id,
            "status": "unknown",
            "awarded_by": None,
            "awarded_to": None,
            "details": [],
            "timeline": [],
            "contract_value": None,
            "extra_points": [],
            "project_intel": {},
            "confidence": 0.0,
            "gaps": [f"Document analysis failed: {exc}"],
            "sources_used": [],
        }

    payload.setdefault("tender_id", tender_id)
    return payload


def summarize_web_findings(
    llm: OllamaChatClient,
    tender_id: str,
    pages: list[dict[str, str]],
    *,
    project_intel: dict[str, Any],
) -> dict[str, Any]:
    compact_pages = [
        {
            "url": page.get("url", ""),
            "title": page.get("title", ""),
            "query": page.get("query", ""),
            "snippet": page.get("snippet", ""),
            "content": (page.get("content", "") or "")[:3000],
            "error": page.get("error", ""),
        }
        for page in pages
    ]

    messages = [
        {
            "role": "system",
            "content": "You are a web research analyst. Return JSON only.",
        },
        {
            "role": "user",
            "content": (
                f"Tender ID: {tender_id}\n"
                f"Project context (use to judge relevance, not only the id):\n"
                f"{json.dumps(project_intel, ensure_ascii=True)[:8000]}\n"
                "From crawled web pages, extract relevant external context and news references.\n"
                "Prefer sources about THIS work / department / geography / NIT reference — ignore unrelated tenders "
                "that merely share a generic id.\n"
                "Return JSON schema:\n"
                "{"
                "\"relevant_news\":[{\"title\":\"...\",\"url\":\"...\",\"summary\":\"...\",\"relevance\":\"high|medium|low\"}],"
                "\"signals\":[\"...\"],"
                "\"risks\":[\"...\"],"
                "\"confidence\":0.0"
                "}\n\n"
                f"Pages:\n{json.dumps(compact_pages, ensure_ascii=True)}"
            ),
        },
    ]

    try:
        output = llm.chat(messages, temperature=0.0, max_tokens=2400, num_ctx=32768)
        payload = extract_json_object(output)
        if not isinstance(payload, dict):
            raise ValueError("web summary returned non-object JSON")
    except Exception as exc:
        payload = {
            "relevant_news": [],
            "signals": [],
            "risks": [f"Web summary failed: {exc}"],
            "confidence": 0.0,
        }
    return payload


def synthesize_final_report(
    llm: OllamaChatClient,
    tender_id: str,
    doc_analysis: dict[str, Any],
    web_summary: dict[str, Any],
    deepresearch: dict[str, Any],
    *,
    project_intel: dict[str, Any],
    max_tokens: int = 6000,
    num_ctx: int = 65536,
) -> dict[str, Any]:
    messages = [
        {
            "role": "system",
            "content": (
                "You are a senior procurement investigator producing an exhaustive tender intelligence report. "
                "You MUST be extremely detailed and cite evidence (quotes + sources). Return JSON only."
            ),
        },
        {
            "role": "user",
            "content": (
                f"Tender ID: {tender_id}\n"
                f"Project / work context (primary identity for this procurement):\n"
                f"{json.dumps(project_intel, ensure_ascii=True)[:16000]}\n"
                "Merge internal tender analysis with deep web/file DeepResearch findings.\n"
                "When describing external research, tie it to THIS work (title, department, geography, NIT/ref) — "
                "not unrelated tenders that share a generic portal id.\n"
                "Return JSON schema:\n"
                "{"
                "\"tender_id\":\"...\","
                "\"executive_summary\":\"...\","
                "\"status\":\"awarded|not_awarded|unknown\","
                "\"awarded_by\":\"...\","
                "\"awarded_to\":\"...\","
                "\"key_points\":[\"...\"],"
                "\"timeline\":[{\"event\":\"...\",\"date\":\"...\"}],"
                "\"web_signals\":[\"...\"],"
                "\"risks\":[\"...\"],"
                "\"red_flags\":[{\"flag\":\"...\",\"why\":\"...\",\"evidence\":[\"...\"]}],"
                "\"news_and_updates\":[{\"title\":\"...\",\"url\":\"...\",\"what_changed\":\"...\",\"date\":\"...\",\"relevance\":\"high|medium|low\"}],"
                "\"file_findings\":[{\"file\":\"...\",\"highlights\":[\"...\"],\"key_quotes\":[\"...\"]}],"
                "\"open_questions\":[\"...\"],"
                "\"sources\":{\"files\":[\"...\"],\"urls\":[\"...\"]},"
                "\"confidence\":0.0"
                "}\n\n"
                f"Document analysis:\n{json.dumps(doc_analysis, ensure_ascii=True)}\n\n"
                f"Web summary:\n{json.dumps(web_summary, ensure_ascii=True)}\n\n"
                f"DeepResearch ledger:\n{json.dumps(deepresearch, ensure_ascii=True)[:220000]}"
            ),
        },
    ]

    try:
        output = llm.chat(
            messages,
            temperature=0.0,
            max_tokens=max_tokens,
            num_ctx=num_ctx,
        )
        payload = extract_json_object(output)
        if not isinstance(payload, dict):
            raise ValueError("final synthesis returned non-object JSON")
    except Exception as exc:
        payload = {
            "tender_id": tender_id,
            "executive_summary": "Unable to synthesize final report.",
            "status": str(doc_analysis.get("status") or "unknown"),
            "awarded_by": str((doc_analysis.get("awarded_by") or {}).get("value") or ""),
            "awarded_to": str((doc_analysis.get("awarded_to") or {}).get("value") or ""),
            "key_points": [],
            "timeline": [],
            "web_signals": [],
            "risks": [f"Final synthesis failed: {exc}"],
            "red_flags": [],
            "news_and_updates": [],
            "file_findings": [],
            "open_questions": [],
            "sources": {"files": [], "urls": []},
            "confidence": 0.0,
        }

    payload.setdefault("tender_id", tender_id)
    return payload


def _dedupe_place_queries(names: list[str], *, limit: int) -> list[str]:
    seen: set[str] = set()
    out: list[str] = []
    for raw in names:
        n = re.sub(r"\s+", " ", (raw or "").strip())[:120]
        if len(n) < 2:
            continue
        key = n.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(n)
        if len(out) >= limit:
            break
    return out


def collect_candidate_place_names(
    project_intel: dict[str, Any],
    doc_analysis: dict[str, Any],
    deepresearch: dict[str, Any],
    *,
    max_candidates: int = 48,
) -> list[str]:
    candidates: list[str] = []

    locs = project_intel.get("locations") or []
    if isinstance(locs, list):
        candidates.extend(str(x).strip() for x in locs if isinstance(x, str) and str(x).strip())

    title = str(project_intel.get("title_or_work") or "")
    for part in re.split(r"[,;]|\bin\b|\bat\b|\bnear\b", title, flags=re.IGNORECASE):
        p = part.strip()
        if len(p) >= 4 and not re.fullmatch(r"\d+[A-Za-z]?", p):
            candidates.append(p)

    details = doc_analysis.get("details") or []
    if isinstance(details, list):
        for line in details:
            if not isinstance(line, str):
                continue
            if re.search(r"\b(location|address|district|mandal|village|city|area|zone)\b", line, re.I):
                idx = line.find(":")
                if idx != -1:
                    tail = line[idx + 1 :].strip()
                    if len(tail) >= 3:
                        candidates.append(tail[:200])

    ev = deepresearch.get("evidence") if isinstance(deepresearch, dict) else None
    ent = (ev or {}).get("entities") if isinstance(ev, dict) else None
    if isinstance(ent, dict):
        locblock = ent.get("locations") or []
        if isinstance(locblock, list):
            candidates.extend(str(x).strip() for x in locblock if isinstance(x, str) and str(x).strip())

    notes = deepresearch.get("research_notes") if isinstance(deepresearch, dict) else None
    if isinstance(notes, list):
        blob = " ".join(str(x) for x in notes[:40] if isinstance(x, str))
        for m in re.finditer(
            r"\b(?:in|at|near)\s+([A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z]+){0,5})\b",
            blob,
        ):
            candidates.append(m.group(1).strip())

    return _dedupe_place_queries(candidates, limit=max_candidates)


def _map_links_for_coordinates(lat: float, lon: float) -> dict[str, str]:
    """
    Prefer stable map URLs for humans. www.openstreetmap.org often returns 429 via Varnish
    when using ?mlat= / ?mlon= query patterns; Google Maps + OSM hash viewer are more reliable.
    """
    return {
        "map_link": f"https://www.google.com/maps?q={lat},{lon}&z=16",
        "google_maps": f"https://www.google.com/maps?q={lat},{lon}&z=16",
        "openstreetmap": f"https://www.openstreetmap.org/#map=16/{lat}/{lon}",
        "geo_uri": f"geo:{lat},{lon}",
    }


def _nominatim_region_suffix(project_intel: dict[str, Any]) -> str:
    title = (str(project_intel.get("title_or_work") or "") + str(project_intel.get("department") or "")).upper()
    loc_blob = " ".join(str(x) for x in (project_intel.get("locations") or []) if isinstance(x, str)).upper()
    if "GHMC" in title or "HYDERABAD" in title or "KUKATPALLY" in title or "KPHB" in loc_blob:
        return ", Hyderabad, Telangana, India"
    if "TELANGANA" in title or "TELANGANA" in loc_blob:
        return ", Telangana, India"
    return ", India"


def _nominatim_get_with_backoff(
    session: requests.Session,
    url: str,
    params: dict[str, Any],
    *,
    timeout_seconds: float,
    max_retries: int = 5,
    base_backoff: float = 2.0,
) -> requests.Response:
    last_resp: requests.Response | None = None
    for attempt in range(max_retries):
        resp = session.get(url, params=params, timeout=timeout_seconds)
        last_resp = resp
        if resp.status_code == 200:
            return resp
        if resp.status_code in (429, 503):
            retry_after = resp.headers.get("Retry-After")
            try:
                wait = float(retry_after) if retry_after else base_backoff * (2**attempt)
            except ValueError:
                wait = base_backoff * (2**attempt)
            wait = min(max(wait, base_backoff), 90.0)
            time.sleep(wait)
            continue
        resp.raise_for_status()
    if last_resp is not None:
        last_resp.raise_for_status()
    raise RuntimeError("Nominatim request failed with no response")


def geocode_places_nominatim(
    place_queries: list[str],
    project_intel: dict[str, Any],
    *,
    delay_seconds: float = 1.1,
    timeout_seconds: float = 25.0,
    limit_per_query: int = 1,
) -> list[dict[str, Any]]:
    """
    Resolve placenames to lat/lon via OpenStreetMap Nominatim (~1 request per second).
    https://nominatim.org/release-docs/develop/api/Search/
    """
    if not place_queries:
        return []

    region_suffix = _nominatim_region_suffix(project_intel)
    session = requests.Session()
    ua = os.getenv(
        "NOMINATIM_USER_AGENT",
        "digital-city-project-tender-research/1.0 (local pipeline)",
    ).strip()
    contact = os.getenv("NOMINATIM_EMAIL", "").strip()
    if contact and contact not in ua:
        ua = f"{ua} contact:{contact}"
    session.headers.update(
        {
            "User-Agent": ua,
            "Accept-Language": "en",
        }
    )

    out: list[dict[str, Any]] = []
    base = "https://nominatim.openstreetmap.org/search"
    last_request_end = 0.0

    for raw in place_queries:
        q = (raw or "").strip()
        if not q:
            continue
        q = q[:200]
        if len(q) < 80 and "india" not in q.lower():
            search_q = f"{q}{region_suffix}"
        else:
            search_q = q

        wait = delay_seconds - (time.monotonic() - last_request_end)
        if wait > 0:
            time.sleep(wait)

        try:
            params: dict[str, Any] = {
                "q": search_q,
                "format": "jsonv2",
                "limit": max(1, limit_per_query),
                "countrycodes": "in",
            }
            r = _nominatim_get_with_backoff(
                session,
                base,
                params,
                timeout_seconds=timeout_seconds,
            )
            rows = r.json()
        except Exception as exc:
            last_request_end = time.monotonic()
            out.append(
                {
                    "query": raw,
                    "search_query": search_q,
                    "lat": None,
                    "lon": None,
                    "display_name": "",
                    "error": str(exc),
                    "map_link": "",
                    "google_maps": "",
                    "openstreetmap": "",
                    "geo_uri": "",
                }
            )
            continue

        last_request_end = time.monotonic()
        if not isinstance(rows, list) or not rows:
            out.append(
                {
                    "query": raw,
                    "search_query": search_q,
                    "lat": None,
                    "lon": None,
                    "display_name": "",
                    "error": "no results",
                    "map_link": "",
                    "google_maps": "",
                    "openstreetmap": "",
                    "geo_uri": "",
                }
            )
            continue

        hit = rows[0]
        try:
            lat_f = float(hit["lat"])
            lon_f = float(hit["lon"])
        except (KeyError, TypeError, ValueError):
            out.append(
                {
                    "query": raw,
                    "search_query": search_q,
                    "lat": None,
                    "lon": None,
                    "display_name": str(hit.get("display_name") or ""),
                    "error": "unparseable coordinates",
                    "map_link": "",
                    "google_maps": "",
                    "openstreetmap": "",
                    "geo_uri": "",
                }
            )
            continue

        display = str(hit.get("display_name") or "")
        links = _map_links_for_coordinates(lat_f, lon_f)
        out.append(
            {
                "query": raw,
                "search_query": search_q,
                "lat": lat_f,
                "lon": lon_f,
                "display_name": display,
                "importance": hit.get("importance"),
                "osm_type": hit.get("osm_type"),
                "osm_id": hit.get("osm_id"),
                "error": "",
                **links,
            }
        )

    return out


def resolve_geotags(
    project_intel: dict[str, Any],
    doc_analysis: dict[str, Any],
    deepresearch: dict[str, Any],
    *,
    skip: bool,
    max_places: int,
    delay_seconds: float,
) -> list[dict[str, Any]]:
    if skip or max_places <= 0:
        return []
    pool = collect_candidate_place_names(
        project_intel,
        doc_analysis,
        deepresearch,
        max_candidates=max(12, max_places * 4),
    )
    return geocode_places_nominatim(pool[:max_places], project_intel, delay_seconds=delay_seconds)


def build_markdown_report_context(
    tender_id: str,
    project_intel: dict[str, Any],
    doc_analysis: dict[str, Any],
    web_summary: dict[str, Any],
    *,
    json_synthesis: dict[str, Any] | None = None,
    geotags: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    """Single bundle passed into the Markdown writer (optional JSON synthesis when --output full)."""
    ctx: dict[str, Any] = {
        "tender_id": tender_id,
        "project_intel": project_intel,
        "document_analysis": doc_analysis,
        "web_summary": web_summary,
    }
    if json_synthesis:
        ctx["json_synthesis"] = json_synthesis
    if geotags is not None:
        ctx["geotags"] = geotags
    return ctx


def fallback_markdown_from_artifacts(
    tender_id: str,
    project_intel: dict[str, Any],
    doc_analysis: dict[str, Any],
    web_summary: dict[str, Any],
    deepresearch: dict[str, Any],
    *,
    reason: str = "",
    geotags: list[dict[str, Any]] | None = None,
) -> str:
    """If the LLM returns nothing, still emit usable Markdown from collected artefacts."""
    parts: list[str] = [
        f"# Tender report: {tender_id}",
        "",
        "## Project context",
        "```json",
        json.dumps(project_intel, indent=2, ensure_ascii=True)[:16000],
        "```",
        "",
        "## Document analysis (raw)",
        "```json",
        json.dumps(doc_analysis, indent=2, ensure_ascii=True)[:24000],
        "```",
        "",
        "## Web summary (raw)",
        "```json",
        json.dumps(web_summary, indent=2, ensure_ascii=True)[:12000],
        "```",
        "",
        "## DeepResearch ledger (raw)",
        "```json",
        json.dumps(deepresearch, indent=2, ensure_ascii=True)[:80000],
        "```",
    ]
    if geotags:
        parts.extend(
            [
                "",
                "## Geotags (Nominatim)",
                "```json",
                json.dumps(geotags, indent=2, ensure_ascii=True)[:24000],
                "```",
            ]
        )
    if reason:
        parts.insert(1, f"> **Note:** Detailed narrative generation did not return text ({reason}). Below is the raw evidence bundle.\n")
    return "\n".join(parts)


def expand_detailed_report_markdown(
    llm: OllamaChatClient,
    *,
    tender_id: str,
    project_intel: dict[str, Any],
    doc_analysis: dict[str, Any],
    web_summary: dict[str, Any],
    deepresearch: dict[str, Any],
    report_context: dict[str, Any],
    geotags: list[dict[str, Any]],
    max_tokens: int,
    num_ctx: int,
) -> str:
    """Long-form Markdown report (plain text output, not JSON)."""
    dr = json.dumps(deepresearch, ensure_ascii=True)
    if len(dr) > 120_000:
        dr = dr[:120_000] + "\n... [truncated deepresearch json]\n"

    ctx = json.dumps(report_context, ensure_ascii=True)
    if len(ctx) > 32000:
        ctx = ctx[:32000] + "\n... [truncated report_context]\n"

    messages = [
        {
            "role": "system",
            "content": (
                "You write exhaustive, well-structured tender intelligence reports in Markdown. "
                "Ground claims in the provided artefacts; quote short evidence snippets; cite file paths or URLs. "
                "Return ONLY Markdown (no JSON, no preamble)."
            ),
        },
        {
            "role": "user",
            "content": (
                f"Tender ID: {tender_id}\n\n"
                "Write a VERY DETAILED Markdown report. Target **at least 4000 words** when the source material "
                "supports it; if facts are genuinely scarce, still produce a thorough, explicit gap analysis.\n\n"
                "Use clear `##` / `###` headings and tables where useful.\n\n"
                "Include these sections in order:\n"
                "1. Executive overview\n"
                "2. Project identity (work title, NIT/ref, department, geography)\n"
                "3. Scope of work & technical description (from internal documents)\n"
                "4. Procurement process (type, evaluation, fees, validity, key conditions)\n"
                "5. Timeline & milestones\n"
                "6. Commercial & financial picture (estimates, bid values if present)\n"
                "7. Award / outcome & bidders (what is known vs unknown)\n"
                "8. External landscape (news, portal pages, related announcements) — **relevance-checked** against "
                "this project, not generic id collisions\n"
                "9. Risks, red flags, conflicts of interest (if any evidence)\n"
                "10. Evidence appendix: bullet list mapping major claims → source path or URL\n"
                "11. **Geography & geotags:** a Markdown table of every resolved place below with columns "
                "`Place (query)` | `Latitude` | `Longitude` | `Label (Nominatim)` | `Google Maps` | `OSM (hash)` | `geo:`. "
                "Use the `google_maps` and `openstreetmap` URLs from the geotags JSON; do **not** invent "
                "`openstreetmap.org/?mlat=` links (they often hit HTTP 429 from Varnish). "
                "Explain briefly how each location relates to the work site or administration. "
                "If a row has no coordinates, say why (ambiguous name, geocoder miss).\n\n"
                f"### Geotags (Nominatim JSON)\n{json.dumps(geotags, ensure_ascii=True)[:12000]}\n\n"
                f"### Project context\n{json.dumps(project_intel, ensure_ascii=True)[:12000]}\n\n"
                f"### Collated context (JSON — may include optional json_synthesis)\n{ctx}\n\n"
                f"### Document analysis (JSON)\n{json.dumps(doc_analysis, ensure_ascii=True)[:24000]}\n\n"
                f"### Web summary (JSON)\n{json.dumps(web_summary, ensure_ascii=True)[:12000]}\n\n"
                f"### DeepResearch ledger (JSON)\n{dr}\n"
            ),
        },
    ]
    return llm.chat(messages, temperature=0.0, max_tokens=max_tokens, num_ctx=num_ctx).strip()


def build_html_report(
    tender_id: str,
    final_summary: dict[str, Any],
    doc_analysis: dict[str, Any],
    web_summary: dict[str, Any],
    web_pages: list[dict[str, str]],
    deepresearch: dict[str, Any],
    issues: list[ExtractionIssue],
    *,
    report_markdown: str = "",
) -> str:
    key_points = final_summary.get("key_points") or []
    timeline = final_summary.get("timeline") or []
    news = web_summary.get("relevant_news") or []
    risks = final_summary.get("risks") or web_summary.get("risks") or []

    def li(items: list[str]) -> str:
        return "".join(f"<li>{pyhtml.escape(str(item))}</li>" for item in items)

    timeline_rows = "".join(
        (
            "<tr>"
            f"<td>{pyhtml.escape(str(item.get('event', '')))}</td>"
            f"<td>{pyhtml.escape(str(item.get('date', '')))}</td>"
            "</tr>"
        )
        for item in timeline
        if isinstance(item, dict)
    )

    news_cards = "".join(
        (
            "<article class='card'>"
            f"<h4>{pyhtml.escape(str(item.get('title', 'Untitled')))}</h4>"
            f"<p>{pyhtml.escape(str(item.get('summary', '')))}</p>"
            f"<p><a href='{pyhtml.escape(str(item.get('url', '#')))}' target='_blank' rel='noopener'>"
            f"{pyhtml.escape(str(item.get('url', '')))}</a></p>"
            "</article>"
        )
        for item in news
        if isinstance(item, dict)
    )

    issue_rows = "".join(
        (
            "<tr>"
            f"<td>{pyhtml.escape(issue.file)}</td>"
            f"<td>{pyhtml.escape(issue.extension)}</td>"
            f"<td>{pyhtml.escape(issue.reason)}</td>"
            "</tr>"
        )
        for issue in issues
    )

    crawled_rows = "".join(
        (
            "<tr>"
            f"<td><a href='{pyhtml.escape(str(page.get('url', '#')))}' target='_blank' rel='noopener'>"
            f"{pyhtml.escape(str(page.get('title', page.get('url', ''))))}</a></td>"
            f"<td>{pyhtml.escape(str(page.get('query', '')))}</td>"
            f"<td>{pyhtml.escape(str(page.get('error', '')))}</td>"
            "</tr>"
        )
        for page in web_pages
    )

    return f"""<!doctype html>
<html lang='en'>
<head>
<meta charset='utf-8'>
<meta name='viewport' content='width=device-width, initial-scale=1'>
<title>Tender Report {pyhtml.escape(tender_id)}</title>
<style>
  :root {{ --bg:#f5f2ea; --ink:#1f2937; --accent:#0f766e; --muted:#6b7280; --card:#ffffff; --line:#d1d5db; }}
  body {{ margin:0; font-family: 'Georgia', 'Times New Roman', serif; background: radial-gradient(circle at top right, #e8f6f3, var(--bg)); color:var(--ink); }}
  .wrap {{ max-width: 1100px; margin: 0 auto; padding: 24px; }}
  h1,h2,h3,h4 {{ margin: 0 0 12px; }}
  h1 {{ font-size: 2rem; color: var(--accent); }}
  .grid {{ display:grid; gap:16px; grid-template-columns: repeat(auto-fit, minmax(300px, 1fr)); }}
  .card {{ background:var(--card); border:1px solid var(--line); border-radius:12px; padding:16px; box-shadow:0 6px 18px rgba(0,0,0,0.05); }}
  .meta {{ color:var(--muted); margin-bottom:12px; }}
  ul {{ margin:0; padding-left:20px; }}
  table {{ width:100%; border-collapse: collapse; background:var(--card); border:1px solid var(--line); border-radius:10px; overflow:hidden; }}
  th, td {{ border-bottom:1px solid var(--line); padding:10px; text-align:left; vertical-align:top; }}
  th {{ background:#ecfeff; }}
  a {{ color:#0c4a6e; }}
  .section {{ margin-top: 20px; }}
  pre {{ white-space: pre-wrap; background:#111827; color:#e5e7eb; padding:12px; border-radius:8px; overflow:auto; }}
</style>
</head>
<body>
  <main class='wrap'>
    <h1>Tender Intelligence Report</h1>
    <p class='meta'>Tender ID: <strong>{pyhtml.escape(tender_id)}</strong> | Generated at: {pyhtml.escape(time.strftime('%Y-%m-%d %H:%M:%S'))}</p>

    <section class='grid'>
      <article class='card'>
        <h2>Executive Summary</h2>
        <p>{pyhtml.escape(str(final_summary.get('executive_summary', '')))}</p>
      </article>
      <article class='card'>
        <h2>Award Snapshot</h2>
        <p><strong>Status:</strong> {pyhtml.escape(str(final_summary.get('status', 'unknown')))}</p>
        <p><strong>Awarded By:</strong> {pyhtml.escape(str(final_summary.get('awarded_by', '')))}</p>
        <p><strong>Awarded To:</strong> {pyhtml.escape(str(final_summary.get('awarded_to', '')))}</p>
        <p><strong>Confidence:</strong> {pyhtml.escape(str(final_summary.get('confidence', 0.0)))}</p>
      </article>
    </section>

    <section class='section card'>
      <h2>Key Points</h2>
      <ul>{li([str(x) for x in key_points])}</ul>
    </section>

    <section class='section card'>
      <h2>Timeline</h2>
      <table>
        <thead><tr><th>Event</th><th>Date</th></tr></thead>
        <tbody>{timeline_rows}</tbody>
      </table>
    </section>

    <section class='section'>
      <h2>News & Web Signals</h2>
      <div class='grid'>{news_cards}</div>
    </section>

    <section class='section card'>
      <h2>Risks</h2>
      <ul>{li([str(x) for x in risks])}</ul>
    </section>

    <section class='section card'>
      <h2>Crawled Sources</h2>
      <table>
        <thead><tr><th>Source</th><th>Query</th><th>Error</th></tr></thead>
        <tbody>{crawled_rows}</tbody>
      </table>
    </section>

    <section class='section card'>
      <h2>Extraction Issues</h2>
      <table>
        <thead><tr><th>File</th><th>Ext</th><th>Issue</th></tr></thead>
        <tbody>{issue_rows}</tbody>
      </table>
    </section>

    <section class='section card'>
      <h2>Raw Objects</h2>
      <h3>DeepResearch</h3>
      <pre>{pyhtml.escape(json.dumps(deepresearch, indent=2, ensure_ascii=True)[:240000])}</pre>
      <h3>Document Analysis</h3>
      <pre>{pyhtml.escape(json.dumps(doc_analysis, indent=2, ensure_ascii=True))}</pre>
      <h3>Web Summary</h3>
      <pre>{pyhtml.escape(json.dumps(web_summary, indent=2, ensure_ascii=True))}</pre>
    </section>

    <section class='section card'>
      <h2>Detailed Report (Markdown)</h2>
      <pre>{pyhtml.escape(str(report_markdown or ''))}</pre>
    </section>
  </main>
</body>
</html>
"""


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Agentic tender pipeline: collect data, research, and write a detailed Markdown report "
            "(default). Use --output full for JSON + HTML + Markdown."
        )
    )
    parser.add_argument("--tender-id", required=True, help="6-digit tender ID")
    parser.add_argument("--downloads-root", default="downloads/agentic_reports")
    parser.add_argument("--output-root", default="outputs/tender_reports")
    parser.add_argument(
        "--output",
        choices=("markdown", "full"),
        default="markdown",
        help="markdown = only {id}_report.md; full = also JSON + HTML.",
    )
    parser.add_argument("--ollama-base-url", default=os.getenv("OLLAMA_BASE_URL", "http://localhost:11434"))
    parser.add_argument("--ollama-model", default=os.getenv("OLLAMA_MODEL", "gemma4:31b-cloud"))
    parser.add_argument("--rate-limit-seconds", type=float, default=2.0)
    parser.add_argument("--llm-timeout-seconds", type=float, default=600.0)
    parser.add_argument("--llm-num-ctx", type=int, default=65536, help="Ollama num_ctx for large prompts.")
    parser.add_argument("--max-doc-files", type=int, default=200)
    parser.add_argument("--max-doc-chars", type=int, default=12000)
    parser.add_argument(
        "--search-results-per-query",
        type=int,
        default=14,
        help="More hits per query → more unique URLs available to crawl.",
    )
    parser.add_argument(
        "--crawl-max-pages",
        type=int,
        default=36,
        help="Max HTML pages to fetch from the initial search-result set.",
    )
    parser.add_argument(
        "--crawl-max-chars",
        type=int,
        default=16000,
        help="Max characters of visible text retained per crawled page.",
    )
    parser.add_argument("--crawl-workers", type=int, default=8)
    parser.add_argument(
        "--max-web-queries",
        type=int,
        default=42,
        help="Cap on distinct web search strings (project-centric queries first).",
    )
    parser.add_argument(
        "--synthesis-max-tokens",
        type=int,
        default=6000,
        help="num_predict for structured JSON synthesis.",
    )
    parser.add_argument(
        "--report-expand-tokens",
        type=int,
        default=12000,
        help="num_predict for long-form Markdown expansion pass.",
    )
    parser.add_argument(
        "--deepresearch-rounds",
        type=int,
        default=4,
        help="How many ReAct research loops to run (web + files) before final write-up.",
    )
    parser.add_argument(
        "--deepresearch-max-pages",
        type=int,
        default=48,
        help="Max pages to crawl during deepresearch (per direct-URL crawl batch).",
    )
    parser.add_argument(
        "--max-geocode-places",
        type=int,
        default=18,
        help="Max distinct placenames to send to Nominatim (0 disables).",
    )
    parser.add_argument(
        "--geocode-delay-seconds",
        type=float,
        default=1.25,
        help="Pause between Nominatim requests (policy: ~1 req/sec; higher reduces HTTP 429).",
    )
    parser.add_argument(
        "--skip-geocode",
        action="store_true",
        help="Do not call OpenStreetMap Nominatim (offline / faster).",
    )
    return parser


def deepresearch_react(
    llm: OllamaChatClient,
    *,
    tender_id: str,
    project_intel: dict[str, Any],
    portal_data: dict[str, Any],
    docs: list[dict[str, str]],
    initial_queries: list[str],
    baseline_pages: list[dict[str, str]],
    rounds: int,
    crawler: LocalNewsCrawler,
    search_results_per_query: int,
    max_pages: int,
) -> dict[str, Any]:
    """
    LangChain-style ReAct orchestration (implemented locally to avoid a brittle
    dependency chain) that alternates between:
      - web_search(query)
      - crawl(urls)
      - read_file(path)
      - search_files(term)
    and accumulates an evidence ledger.
    """

    def tool_web_search(query: str) -> list[dict[str, str]]:
        results = crawler.search([query], max_results_per_query=search_results_per_query)
        return [asdict(r) for r in results]

    def tool_web_crawl(urls: list[str]) -> list[dict[str, str]]:
        items = [SearchResult(title=u, href=u, snippet="", source_query="direct_url") for u in urls if u]
        return crawler.crawl(items, max_pages=min(len(items), max_pages))

    def tool_file_manifest() -> list[dict[str, Any]]:
        return [{"path": d["path"], "chars": len(d.get("text") or "")} for d in docs]

    def tool_read_file(path: str) -> dict[str, str]:
        for d in docs:
            if d.get("path") == path:
                return {"path": path, "text": (d.get("text") or "")[:12000]}
        return {"path": path, "text": ""}

    def tool_search_files(term: str) -> list[dict[str, str]]:
        if not term:
            return []
        needle = term.lower()
        hits: list[dict[str, str]] = []
        for d in docs:
            text = (d.get("text") or "")
            if needle in text.lower():
                idx = text.lower().find(needle)
                start = max(0, idx - 400)
                end = min(len(text), idx + 1200)
                hits.append({"path": d.get("path", ""), "snippet": text[start:end]})
        return hits[:25]

    ledger: dict[str, Any] = {
        "tender_id": tender_id,
        "project_intel": project_intel,
        "rounds": rounds,
        "research_notes": [],
        "web_pages": list(baseline_pages),
        "queries_run": list(initial_queries),
        "files_manifest": tool_file_manifest(),
        "evidence": {
            "file_quotes": [],
            "web_quotes": [],
            "entities": {"people": [], "orgs": [], "locations": [], "products": [], "ids": []},
            "dates": [],
            "money": [],
            "claims": [],
        },
    }

    system = (
        "You are DeepResearch, an investigative tender research agent. "
        "You must uncover everything relevant about this tender, prioritising:\n"
        "- news, updates, corrigenda, award status, litigation/complaints\n"
        "- contracting authority, vendors, contract value, scope, milestones\n"
        "- evaluation criteria, eligibility, risks, red flags, conflicts\n"
        "Search the WEB using the **project/work identity** (title, department, NIT/ref, geography). "
        "Do **not** issue queries that are only the 6-digit portal tender id — combine the id with "
        "department / work title / place names when needed.\n"
        "You can use tools. When you need evidence, call tools instead of guessing.\n"
        "Return STRICT JSON only."
    )

    def ask_for_plan(existing_notes: list[str]) -> dict[str, Any]:
        messages = [
            {"role": "system", "content": system},
            {
                "role": "user",
                "content": (
                    f"Tender ID: {tender_id}\n"
                    f"Project / work context (use for web_search queries):\n"
                    f"{json.dumps(project_intel, ensure_ascii=True)[:12000]}\n\n"
                    "You have the following tools available:\n"
                    "- web_search(query) -> list of {title,href,snippet,source_query}\n"
                    "- web_crawl(urls[]) -> list of {url,title,query,snippet,content,error}\n"
                    "- file_manifest() -> list of {path,chars}\n"
                    "- read_file(path) -> {path,text}\n"
                    "- search_files(term) -> list of {path,snippet}\n\n"
                    "Create a ReAct research step for THIS round. "
                    "Decide what to search next, which files to inspect, and what to extract.\n"
                    "Return JSON:\n"
                    "{"
                    '"thoughts":"brief",'
                    '"web_queries":["..."],'
                    '"urls_to_crawl":["..."],'
                    '"file_paths_to_read":["..."],'
                    '"file_terms_to_search":["..."],'
                    '"extraction_targets":["specific facts to extract"]'
                    "}\n\n"
                    f"Portal data (compact): {json.dumps(portal_data, ensure_ascii=True)[:14000]}\n\n"
                    f"Existing notes: {json.dumps(existing_notes, ensure_ascii=True)[:8000]}\n"
                ),
            },
        ]
        raw = llm.chat(messages, temperature=0.0, max_tokens=1400)
        return extract_json_object(raw)

    def ask_for_extraction(
        *,
        round_idx: int,
        plan: dict[str, Any],
        file_reads: list[dict[str, str]],
        file_hits: list[dict[str, str]],
        web_results: list[dict[str, Any]],
        crawled_pages: list[dict[str, str]],
    ) -> dict[str, Any]:
        messages = [
            {"role": "system", "content": system},
            {
                "role": "user",
                "content": (
                    f"Tender ID: {tender_id}\nRound: {round_idx}\n"
                    f"Project / work context:\n{json.dumps(project_intel, ensure_ascii=True)[:8000]}\n\n"
                    "Extract VERIFIED facts only, each with at least one quote and a source.\n"
                    "Return JSON:\n"
                    "{"
                    '"notes":["..."],'
                    '"file_quotes":[{"path":"...","quote":"...","why":"..."}],'
                    '"web_quotes":[{"url":"...","quote":"...","why":"..."}],'
                    '"entities":{"people":["..."],"orgs":["..."],"locations":["..."],"products":["..."],"ids":["..."]},'
                    '"dates":[{"date":"...","context":"...","source":"file|web","ref":"path-or-url"}],'
                    '"money":[{"amount":"...","context":"...","source":"file|web","ref":"path-or-url"}],'
                    '"claims":[{"claim":"...","supporting_refs":["..."],"confidence":"high|medium|low"}],'
                    '"next_queries":["..."]'
                    "}\n\n"
                    f"Plan: {json.dumps(plan, ensure_ascii=True)[:8000]}\n\n"
                    f"Web search results: {json.dumps(web_results, ensure_ascii=True)[:12000]}\n\n"
                    f"Crawled pages: {json.dumps(crawled_pages, ensure_ascii=True)[:18000]}\n\n"
                    f"File term hits: {json.dumps(file_hits, ensure_ascii=True)[:12000]}\n\n"
                    f"File reads: {json.dumps(file_reads, ensure_ascii=True)[:18000]}\n"
                ),
            },
        ]
        raw = llm.chat(messages, temperature=0.0, max_tokens=2200)
        return extract_json_object(raw)

    notes: list[str] = []
    for r in range(1, max(1, rounds) + 1):
        plan = ask_for_plan(notes)

        web_queries = [str(q) for q in (plan.get("web_queries") or []) if isinstance(q, str) and q.strip()]
        urls_to_crawl = [str(u) for u in (plan.get("urls_to_crawl") or []) if isinstance(u, str) and u.strip()]
        file_paths_to_read = [str(p) for p in (plan.get("file_paths_to_read") or []) if isinstance(p, str) and p.strip()]
        file_terms_to_search = [str(t) for t in (plan.get("file_terms_to_search") or []) if isinstance(t, str) and t.strip()]

        web_results: list[dict[str, Any]] = []
        for q in web_queries[:10]:
            ledger["queries_run"].append(q)
            web_results.extend(tool_web_search(q))

        crawled = tool_web_crawl(urls_to_crawl[: max_pages])
        if crawled:
            ledger["web_pages"].extend(crawled)

        file_reads = [tool_read_file(p) for p in file_paths_to_read[:14]]
        file_hits: list[dict[str, str]] = []
        for term in file_terms_to_search[:10]:
            file_hits.extend(tool_search_files(term))

        extracted = ask_for_extraction(
            round_idx=r,
            plan=plan,
            file_reads=file_reads,
            file_hits=file_hits,
            web_results=web_results,
            crawled_pages=crawled,
        )

        new_notes = [str(x) for x in (extracted.get("notes") or []) if isinstance(x, str) and x.strip()]
        notes.extend(new_notes)
        ledger["research_notes"].extend(new_notes)

        for key in ("file_quotes", "web_quotes"):
            items = extracted.get(key) or []
            if isinstance(items, list):
                ledger["evidence"][key].extend([x for x in items if isinstance(x, dict)])

        ent = extracted.get("entities") or {}
        if isinstance(ent, dict):
            for k in ("people", "orgs", "locations", "products", "ids"):
                vals = ent.get(k) or []
                if isinstance(vals, list):
                    ledger["evidence"]["entities"][k].extend([str(v) for v in vals if isinstance(v, str) and v.strip()])

        for key in ("dates", "money", "claims"):
            items = extracted.get(key) or []
            if isinstance(items, list):
                ledger["evidence"][key].extend([x for x in items if isinstance(x, dict)])

        next_queries = [str(q) for q in (extracted.get("next_queries") or []) if isinstance(q, str) and q.strip()]
        # Seed the next round with new queries by appending to the note list.
        if next_queries:
            notes.append("Next round query ideas: " + "; ".join(next_queries[:10]))

    # Deduplicate entity lists
    for k in ("people", "orgs", "locations", "products", "ids"):
        seen: set[str] = set()
        uniq: list[str] = []
        for v in ledger["evidence"]["entities"][k]:
            vv = v.strip()
            if vv and vv not in seen:
                seen.add(vv)
                uniq.append(vv)
        ledger["evidence"]["entities"][k] = uniq

    return ledger


def deepresearch_langchain_react(
    *,
    ollama_base_url: str,
    ollama_model: str,
    tender_id: str,
    project_intel: dict[str, Any],
    portal_data: dict[str, Any],
    docs: list[dict[str, str]],
    initial_queries: list[str],
    baseline_pages: list[dict[str, str]],
    rounds: int,
    crawler: LocalNewsCrawler,
    search_results_per_query: int,
    max_pages: int,
) -> dict[str, Any]:
    """
    Preferred implementation: LangChain ReAct agent with Ollama LLM.
    If LangChain isn't available at runtime, callers should fall back to `deepresearch_react`.
    """
    from langchain.agents import AgentExecutor, create_react_agent  # type: ignore
    from langchain_core.prompts import PromptTemplate  # type: ignore
    from langchain_core.tools import tool  # type: ignore
    from langchain_ollama import ChatOllama  # type: ignore

    llm = ChatOllama(model=ollama_model, base_url=ollama_base_url.rstrip("/"), temperature=0.0)

    @tool
    def web_search(query: str) -> str:
        "Search the web for relevant pages; returns JSON list of results."
        results = crawler.search([query], max_results_per_query=search_results_per_query)
        return json.dumps([asdict(r) for r in results], ensure_ascii=True)

    @tool
    def web_crawl(urls_json: str) -> str:
        "Crawl the given URLs (JSON list of strings) and return JSON list of page objects."
        try:
            urls = json.loads(urls_json)
        except Exception:
            urls = []
        if not isinstance(urls, list):
            urls = []
        items = [
            SearchResult(title=str(u), href=str(u), snippet="", source_query="direct_url")
            for u in urls
            if isinstance(u, str) and u.strip()
        ]
        pages = crawler.crawl(items, max_pages=min(len(items), max_pages))
        return json.dumps(pages, ensure_ascii=True)

    @tool
    def file_manifest(_: str = "") -> str:
        "Return JSON list of all available file paths and char counts."
        manifest = [{"path": d["path"], "chars": len(d.get("text") or "")} for d in docs]
        return json.dumps(manifest, ensure_ascii=True)

    @tool
    def read_file(path: str) -> str:
        "Read a file by exact path from the tender corpus and return JSON {path,text}."
        for d in docs:
            if d.get("path") == path:
                payload = {"path": path, "text": (d.get("text") or "")[:12000]}
                return json.dumps(payload, ensure_ascii=True)
        return json.dumps({"path": path, "text": ""}, ensure_ascii=True)

    @tool
    def search_files(term: str) -> str:
        "Search across file texts for a term; return JSON list of {path,snippet}."
        term = (term or "").strip()
        if not term:
            return "[]"
        needle = term.lower()
        hits: list[dict[str, str]] = []
        for d in docs:
            text = (d.get("text") or "")
            lower = text.lower()
            if needle in lower:
                idx = lower.find(needle)
                start = max(0, idx - 400)
                end = min(len(text), idx + 1200)
                hits.append({"path": d.get("path", ""), "snippet": text[start:end]})
        return json.dumps(hits[:25], ensure_ascii=True)

    tools = [web_search, web_crawl, file_manifest, read_file, search_files]

    prompt = PromptTemplate.from_template(
        "You are DeepResearch, an investigative tender research agent.\n"
        "Your job: find everything possible about the tender, focusing on news, updates, award status, "
        "corrigenda, vendors, contract value, scope, evaluation, timelines, risks.\n"
        "Use web_search with **project/work details** (title, department, NIT/ref, geography). "
        "Avoid queries that are ONLY the 6-digit portal tender id.\n"
        "You MUST use tools to gather evidence, and you MUST cite sources in your final JSON.\n\n"
        "Tender ID: {tender_id}\n"
        "Project / work context: {project_intel}\n"
        "Portal data (compact): {portal_data}\n"
        "Initial queries: {initial_queries}\n\n"
        "Use the following format:\n\n"
        "Question: the question you must answer\n"
        "Thought: you should always think about what to do\n"
        "Action: the action to take, should be one of [{tool_names}]\n"
        "Action Input: the input to the action\n"
        "Observation: the result of the action\n"
        "... (repeat Thought/Action/Action Input/Observation as needed)\n"
        "Final Answer: return STRICT JSON with keys:\n"
        "{"
        "\"tender_id\":\"...\","
        "\"research_notes\":[\"...\"],"
        "\"queries_run\":[\"...\"],"
        "\"web_pages\":[{\"url\":\"...\",\"title\":\"...\",\"query\":\"...\",\"snippet\":\"...\",\"content\":\"...\",\"error\":\"...\"}],"
        "\"evidence\":{"
        "\"file_quotes\":[{\"path\":\"...\",\"quote\":\"...\",\"why\":\"...\"}],"
        "\"web_quotes\":[{\"url\":\"...\",\"quote\":\"...\",\"why\":\"...\"}],"
        "\"entities\":{\"people\":[],\"orgs\":[],\"locations\":[],\"products\":[],\"ids\":[]},"
        "\"dates\":[],\"money\":[],\"claims\":[]"
        "}"
        "}\n\n"
        "Begin.\n"
        "Question: Build an exhaustive evidence ledger for this tender.\n"
        "{agent_scratchpad}"
    )

    agent = create_react_agent(llm, tools, prompt)
    executor = AgentExecutor(agent=agent, tools=tools, verbose=False, handle_parsing_errors=True, max_iterations=max(6, rounds * 4))

    result = executor.invoke(
        {
            "tender_id": tender_id,
            "project_intel": json.dumps(project_intel, ensure_ascii=True)[:12000],
            "portal_data": json.dumps(portal_data, ensure_ascii=True)[:14000],
            "initial_queries": json.dumps(initial_queries, ensure_ascii=True)[:4000],
        }
    )
    output_text = result.get("output", "") if isinstance(result, dict) else str(result)
    payload = extract_json_object(output_text)
    if not isinstance(payload, dict):
        raise ValueError("LangChain ReAct returned non-object JSON")

    # merge baseline pages (so we preserve the earlier crawl)
    pages = payload.get("web_pages")
    if isinstance(pages, list):
        merged = list(baseline_pages) + [p for p in pages if isinstance(p, dict)]
        payload["web_pages"] = merged
    else:
        payload["web_pages"] = list(baseline_pages)

    # guarantee a couple of expected keys
    payload.setdefault("tender_id", tender_id)
    payload.setdefault("queries_run", list(initial_queries))
    payload.setdefault("project_intel", project_intel)
    return payload


def main() -> int:
    args = build_parser().parse_args()
    tender_id = args.tender_id.strip()
    if not re.fullmatch(r"\d{6}", tender_id):
        print("tender-id must be exactly 6 digits")
        return 2

    downloads_root = Path(args.downloads_root).resolve()
    output_root = Path(args.output_root).resolve()
    tender_output_dir = output_root / tender_id
    tender_output_dir.mkdir(parents=True, exist_ok=True)

    print(f"[1/6] Collecting tender data for {tender_id}")
    collected = collect_tender_data(tender_id, downloads_root, args.rate_limit_seconds)

    source_dirs = [Path(p) for p in collected["source_dirs"]]
    print(f"[2/6] Reading documents from {len(source_dirs)} source directories")
    reader = DocumentReader(max_files=args.max_doc_files, max_chars_per_file=args.max_doc_chars)
    docs, issues = reader.read_documents(source_dirs)
    print(f"      Readable docs: {len(docs)} | Issues: {len(issues)}")

    llm = OllamaChatClient(
        base_url=args.ollama_base_url,
        model=args.ollama_model,
        timeout_seconds=args.llm_timeout_seconds,
    )

    print("[3/6] Analyzing tender documents with LLM")
    doc_analysis = analyze_documents(llm, tender_id, docs, collected["portal_data"])
    project_intel = merge_project_intel(doc_analysis)

    queries = build_web_queries(
        tender_id,
        doc_analysis,
        project_intel=project_intel,
        max_queries=max(12, args.max_web_queries),
    )
    print(f"[4/6] Web research with local crawler. Queries ({len(queries)}): {queries[:8]}{' ...' if len(queries) > 8 else ''}")
    crawler = LocalNewsCrawler(max_workers=args.crawl_workers)
    search_results = crawler.search(queries, max_results_per_query=args.search_results_per_query)
    crawled_pages = crawler.crawl(
        search_results,
        max_pages=max(1, args.crawl_max_pages),
        max_chars=max(2000, args.crawl_max_chars),
    )
    print(f"      Search results: {len(search_results)} | Crawled pages: {len(crawled_pages)}")

    print("[5/6] DeepResearch (ReAct loop) + Markdown report")
    deepresearch: dict[str, Any]
    try:
        deepresearch = deepresearch_langchain_react(
            ollama_base_url=args.ollama_base_url,
            ollama_model=args.ollama_model,
            tender_id=tender_id,
            project_intel=project_intel,
            portal_data=collected["portal_data"],
            docs=docs,
            initial_queries=queries,
            baseline_pages=crawled_pages,
            rounds=args.deepresearch_rounds,
            crawler=crawler,
            search_results_per_query=args.search_results_per_query,
            max_pages=args.deepresearch_max_pages,
        )
    except Exception:
        deepresearch = deepresearch_react(
            llm,
            tender_id=tender_id,
            project_intel=project_intel,
            portal_data=collected["portal_data"],
            docs=docs,
            initial_queries=queries,
            baseline_pages=crawled_pages,
            rounds=args.deepresearch_rounds,
            crawler=crawler,
            search_results_per_query=args.search_results_per_query,
            max_pages=args.deepresearch_max_pages,
        )

    web_summary = summarize_web_findings(
        llm,
        tender_id,
        deepresearch.get("web_pages") or crawled_pages,
        project_intel=project_intel,
    )

    geotags = resolve_geotags(
        project_intel,
        doc_analysis,
        deepresearch,
        skip=args.skip_geocode,
        max_places=max(0, args.max_geocode_places),
        delay_seconds=max(0.5, args.geocode_delay_seconds),
    )
    print(f"      Geotags resolved: {sum(1 for g in geotags if g.get('lat') is not None)}/{len(geotags)}")

    json_synthesis: dict[str, Any] | None = None
    if args.output == "full":
        json_synthesis = synthesize_final_report(
            llm,
            tender_id,
            doc_analysis,
            web_summary,
            deepresearch,
            project_intel=project_intel,
            max_tokens=max(1024, args.synthesis_max_tokens),
            num_ctx=max(8192, args.llm_num_ctx),
        )

    report_context = build_markdown_report_context(
        tender_id,
        project_intel,
        doc_analysis,
        web_summary,
        json_synthesis=json_synthesis,
        geotags=geotags,
    )

    md_err = ""
    try:
        report_md = expand_detailed_report_markdown(
            llm,
            tender_id=tender_id,
            project_intel=project_intel,
            doc_analysis=doc_analysis,
            web_summary=web_summary,
            deepresearch=deepresearch,
            report_context=report_context,
            geotags=geotags,
            max_tokens=max(2048, args.report_expand_tokens),
            num_ctx=max(8192, args.llm_num_ctx),
        )
    except Exception as exc:
        md_err = str(exc)
        report_md = ""

    if not (report_md or "").strip():
        report_md = fallback_markdown_from_artifacts(
            tender_id,
            project_intel,
            doc_analysis,
            web_summary,
            deepresearch,
            reason=md_err or "empty model response",
            geotags=geotags,
        )

    md_path = tender_output_dir / f"{tender_id}_report.md"
    print("[6/6] Writing report")
    write_text(md_path, report_md)
    print(f"Markdown report: {md_path}")

    if args.output != "full":
        return 0

    final_summary = json_synthesis or {
        "tender_id": tender_id,
        "executive_summary": "",
        "status": str(doc_analysis.get("status") or "unknown"),
        "key_points": [],
        "timeline": [],
        "web_signals": [],
        "risks": [],
        "red_flags": [],
        "news_and_updates": [],
        "file_findings": [],
        "open_questions": [],
        "sources": {"files": [], "urls": []},
        "confidence": 0.0,
    }

    payload = {
        "tender_id": tender_id,
        "generated_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
        "inputs": {
            "downloads_root": str(downloads_root),
            "output_root": str(output_root),
            "ollama_model": args.ollama_model,
            "ollama_base_url": args.ollama_base_url,
            "output": args.output,
            "deepresearch_rounds": args.deepresearch_rounds,
            "deepresearch_max_pages": args.deepresearch_max_pages,
            "max_web_queries": args.max_web_queries,
            "synthesis_max_tokens": args.synthesis_max_tokens,
            "report_expand_tokens": args.report_expand_tokens,
            "llm_num_ctx": args.llm_num_ctx,
            "crawl_max_pages": args.crawl_max_pages,
            "crawl_max_chars": args.crawl_max_chars,
            "search_results_per_query": args.search_results_per_query,
            "max_geocode_places": args.max_geocode_places,
            "geocode_delay_seconds": args.geocode_delay_seconds,
            "skip_geocode": args.skip_geocode,
        },
        "geotags": geotags,
        "collection": collected,
        "document_count": len(docs),
        "extraction_issues": [asdict(issue) for issue in issues],
        "project_intel": project_intel,
        "document_analysis": doc_analysis,
        "web_queries": queries,
        "web_search_results": [asdict(item) for item in search_results],
        "web_crawled_pages": crawled_pages,
        "deepresearch": deepresearch,
        "web_summary": web_summary,
        "final_summary": final_summary,
        "detailed_report_markdown": report_md,
    }

    json_path = tender_output_dir / f"{tender_id}_report.json"
    html_path = tender_output_dir / f"{tender_id}_report.html"
    write_json(json_path, payload)
    html_text = build_html_report(
        tender_id,
        final_summary,
        doc_analysis,
        web_summary,
        crawled_pages,
        deepresearch,
        issues,
        report_markdown=report_md,
    )
    write_text(html_path, html_text)
    print(f"JSON report: {json_path}")
    print(f"HTML report: {html_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
