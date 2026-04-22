"""
tender_pipeline.py
==================
Agentic pipeline that extracts DEEPLY DETAILED structured tender award data
from procurement documents.

Supported backends
------------------
  ollama  (default) – native /api/chat, any model pulled via `ollama pull`
  vllm              – OpenAI-compatible /chat/completions (openai SDK)

Key improvements over v1
------------------------
* Reads EVERY file type: .docx/.docm, .pdf, .doc, .xlsx/.xlsm/.xls,
  .pptx/.ppt, .odt/.ods, .rtf, .txt/.csv/.tsv/.md/.json/.xml/.html/.log,
  .yaml/.yml – nothing is skipped unless it is a true binary blob.
* Reads the tender root itself and every sub-directory, not just *_extracted.
* Automatically unpacks .zip/.tar/.gz archives encountered while walking.
* Deep extraction prompts demanding full financial breakdowns, all parties,
  all dates, all clauses, evaluation criteria, scoring, bid values, penalties,
  SLA targets, amendment history, delivery schedules, contacts, legal refs.
* Output JSON is exhaustive – every field populated, no skeleton placeholders.
* Heartbeat + --limit + --verbose from the reference implementation.
* Thread-safe JSONL appending.

Usage
-----
  # Ollama (default):
  python tender_pipeline.py --ollama-model mistral

  # vLLM:
  python tender_pipeline.py --backend vllm \\
      --vllm-base-url http://10.4.25.56:8000/v1 \\
      --vllm-model Qwen/Qwen3-30B-A3B-GPTQ-Int4

Env vars (all optional):
  OLLAMA_BASE_URL   default: http://localhost:11434
  OLLAMA_MODEL      default: llama3
  VLLM_BASE_URL     default: http://localhost:8000/v1
  VLLM_MODEL        default: Qwen/Qwen3-30B-A3B-GPTQ-Int4
  VLLM_API_KEY      default: EMPTY
"""

from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import threading
import time
import zipfile
from abc import ABC, abstractmethod
from concurrent.futures import FIRST_COMPLETED, ThreadPoolExecutor, wait
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

import requests
from lxml import etree, html

# ---------------------------------------------------------------------------
# File-type sets
# ---------------------------------------------------------------------------

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".tsv", ".json", ".xml",
    ".html", ".htm", ".log", ".yaml", ".yml", ".rtf",
    ".ini", ".cfg", ".conf", ".properties", ".toml",
    ".sql", ".tex", ".rst",
}

# Only truly unreadable binary blobs.
SKIP_BINARY_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".ico",
    ".mp3", ".mp4", ".avi", ".mov", ".wav",
    ".exe", ".dll", ".so", ".dylib",
    ".bin", ".dat", ".db", ".sqlite",
    ".iso", ".img",
}

# Archives we unpack recursively.
ARCHIVE_EXTENSIONS = {".zip", ".tar", ".gz", ".bz2", ".7z", ".rar"}


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass(slots=True)
class ExtractionIssue:
    file: str
    extension: str
    reason: str


@dataclass(slots=True)
class FileDocument:
    path: Path
    rel_path: str
    text: str
    size_bytes: int
    extension: str


# ---------------------------------------------------------------------------
# LLM abstraction
# ---------------------------------------------------------------------------

class BaseLLMClient(ABC):
    @abstractmethod
    def chat(
        self,
        messages: list[dict[str, str]],
        *,
        temperature: float = 0.0,
        max_tokens: int = 2000,
    ) -> str: ...


class OllamaClient(BaseLLMClient):
    """
    Native Ollama /api/chat client.
    Docs: https://github.com/ollama/ollama/blob/main/docs/api.md
    """

    def __init__(
        self,
        base_url: str = "http://localhost:11434",
        model: str = "llama3",
        timeout_seconds: float = 300.0,
    ) -> None:
        self.base_url = base_url.rstrip("/")
        self.model = model
        self.timeout_seconds = timeout_seconds
        self._local = threading.local()

    def _session(self) -> requests.Session:
        if not hasattr(self._local, "session"):
            self._local.session = requests.Session()
        return self._local.session

    def chat(
        self,
        messages: list[dict[str, str]],
        *,
        temperature: float = 0.0,
        max_tokens: int = 2000,
    ) -> str:
        payload: dict[str, Any] = {
            "model": self.model,
            "messages": messages,
            "stream": False,
            "options": {
                "temperature": temperature,
                "num_predict": max_tokens,
                "num_ctx": 16384,
            },
        }
        r = self._session().post(
            f"{self.base_url}/api/chat",
            json=payload,
            headers={"Content-Type": "application/json"},
            timeout=self.timeout_seconds,
        )
        r.raise_for_status()
        return r.json()["message"]["content"]

    def available_models(self) -> list[str]:
        try:
            r = self._session().get(f"{self.base_url}/api/tags", timeout=10)
            r.raise_for_status()
            return [m["name"] for m in r.json().get("models", [])]
        except Exception:
            return []

    def check_model(self) -> bool:
        return any(
            n == self.model or n.startswith(f"{self.model}:")
            for n in self.available_models()
        )


class VllmChatClient(BaseLLMClient):
    """OpenAI-compatible endpoint (vLLM, LM Studio, llama.cpp …)."""

    def __init__(
        self,
        base_url: str,
        model: str,
        api_key: str = "EMPTY",
        timeout_seconds: float = 180.0,
    ) -> None:
        from openai import OpenAI  # type: ignore
        self.model = model
        self.client = OpenAI(
            base_url=base_url.rstrip("/"),
            api_key=api_key or "EMPTY",
            timeout=timeout_seconds,
        )

    def chat(
        self,
        messages: list[dict[str, str]],
        *,
        temperature: float = 0.0,
        max_tokens: int = 2000,
    ) -> str:
        resp = self.client.chat.completions.create(
            model=self.model,
            messages=messages,
            temperature=temperature,
            max_tokens=max_tokens,
        )
        content = resp.choices[0].message.content
        if content is None:
            raise ValueError("Model returned empty content")
        return content


# ---------------------------------------------------------------------------
# Pipeline
# ---------------------------------------------------------------------------

class TenderAgenticPipeline:
    def __init__(
        self,
        llm: BaseLLMClient,
        *,
        max_files: int = 300,
        max_file_chars: int = 18_000,
        max_relevant_files: int = 20,
        unpack_archives: bool = True,
        verbose: bool = False,
    ) -> None:
        self.llm = llm
        self.max_files = max_files
        self.max_file_chars = max_file_chars
        self.max_relevant_files = max_relevant_files
        self.unpack_archives = unpack_archives
        self.verbose = verbose
        self._write_lock = threading.Lock()

    # ------------------------------------------------------------------
    # Public entry
    # ------------------------------------------------------------------

    def process_tender(self, tender_id: str, tender_root: Path) -> dict[str, Any]:
        self._log(f"[{tender_id}] Scanning {tender_root}")
        docs, issues = self._read_all_documents(tender_root)
        self._log(f"[{tender_id}] Read {len(docs)} documents, {len(issues)} issues")

        if not docs:
            return {
                "tender_id": tender_id,
                "status": "unknown",
                "error": "No readable files found.",
                "tender_root": str(tender_root),
                "extraction_issues": [asdict(i) for i in issues],
            }

        relevant = self._triage_relevant_files(tender_id, docs)
        self._log(f"[{tender_id}] Selected {len(relevant)} relevant files: {relevant}")

        evidence = self._extract_deep_evidence(tender_id, docs, relevant)
        result = self._synthesize_detailed(tender_id, docs, evidence)
        result["extraction_issues"] = [asdict(i) for i in issues]
        result["all_files_found"] = [
            {
                "path": d.rel_path,
                "extension": d.extension,
                "size_bytes": d.size_bytes,
                "chars_extracted": len(d.text),
            }
            for d in docs
        ]
        return result

    def _log(self, msg: str) -> None:
        if self.verbose:
            print(msg)

    # ------------------------------------------------------------------
    # Document discovery & reading
    # ------------------------------------------------------------------

    def _read_all_documents(
        self, root: Path
    ) -> tuple[list[FileDocument], list[ExtractionIssue]]:
        import tempfile

        docs: list[FileDocument] = []
        issues: list[ExtractionIssue] = []
        tmp_dirs: list[Path] = []
        all_files: list[tuple[Path, Path]] = []  # (file_path, base_root)

        def collect(directory: Path, base: Path) -> None:
            for p in sorted(directory.rglob("*"), key=lambda x: str(x).lower()):
                if not p.is_file():
                    continue
                suf = p.suffix.lower()
                if suf in ARCHIVE_EXTENSIONS and self.unpack_archives:
                    tmp = Path(tempfile.mkdtemp(prefix="tender_unpack_"))
                    tmp_dirs.append(tmp)
                    if self._unpack_archive(p, tmp):
                        collect(tmp, tmp)
                    else:
                        try:
                            rel = str(p.relative_to(base)).replace("\\", "/")
                        except ValueError:
                            rel = p.name
                        issues.append(ExtractionIssue(
                            file=rel, extension=suf, reason="Archive unpack failed"
                        ))
                else:
                    all_files.append((p, base))

        collect(root, root)

        # Deduplicate by normalised relative path
        seen: set[str] = set()
        unique: list[tuple[Path, Path]] = []
        for fp, base in all_files:
            try:
                key = str(fp.relative_to(base)).replace("\\", "/").lower()
            except ValueError:
                key = str(fp).lower()
            if key not in seen:
                seen.add(key)
                unique.append((fp, base))

        unique = unique[: self.max_files]

        for path, base in unique:
            text, handled, reason = self._extract_text(path)
            text = self._normalize_text(text)

            try:
                rel_path = str(path.relative_to(base)).replace("\\", "/")
            except ValueError:
                rel_path = path.name

            if reason:
                issues.append(ExtractionIssue(
                    file=rel_path,
                    extension=path.suffix.lower(),
                    reason=reason,
                ))

            if not handled or not text:
                continue

            docs.append(FileDocument(
                path=path,
                rel_path=rel_path,
                text=text[: self.max_file_chars],
                size_bytes=path.stat().st_size,
                extension=path.suffix.lower(),
            ))

        for tmp in tmp_dirs:
            shutil.rmtree(tmp, ignore_errors=True)

        return docs, issues

    # ------------------------------------------------------------------
    # Archive unpacking
    # ------------------------------------------------------------------

    @staticmethod
    def _unpack_archive(path: Path, dest: Path) -> bool:
        suf = path.suffix.lower()
        try:
            if suf == ".zip":
                with zipfile.ZipFile(path) as zf:
                    zf.extractall(dest)
                return True
            if suf in {".tar", ".gz", ".bz2"}:
                import tarfile
                with tarfile.open(path) as tf:
                    tf.extractall(dest)
                return True
            if suf in {".7z", ".rar"}:
                exe = shutil.which("7z") or shutil.which("7za")
                if exe:
                    r = subprocess.run(
                        [exe, "x", str(path), f"-o{dest}", "-y"],
                        capture_output=True, timeout=60,
                    )
                    return r.returncode == 0
        except Exception:
            pass
        return False

    # ------------------------------------------------------------------
    # Text extraction
    # ------------------------------------------------------------------

    def _extract_text(self, path: Path) -> tuple[str, bool, str | None]:
        suf = path.suffix.lower()

        if suf in SKIP_BINARY_EXTENSIONS:
            return "", True, None

        if suf in TEXT_EXTENSIONS:
            for enc in ("utf-8", "utf-16", "latin-1", "cp1252"):
                try:
                    raw = path.read_text(encoding=enc, errors="ignore")
                    if raw.strip():
                        if suf == ".rtf":
                            return self._extract_rtf(raw), True, None
                        return raw, True, None
                except OSError:
                    break
            return "", True, f"Could not read text file: {path.name}"

        dispatch: dict[str, Any] = {
            ".docx": self._extract_docx,
            ".docm": self._extract_docx,
            ".pdf":  self._extract_pdf,
            ".doc":  self._extract_doc,
            ".xlsx": self._extract_xlsx,
            ".xlsm": self._extract_xlsx,
            ".xls":  self._extract_xls,
            ".pptx": self._extract_pptx,
            ".ppt":  self._extract_ppt_legacy,
            ".odt":  self._extract_opendoc,
            ".ods":  self._extract_opendoc,
            ".odp":  self._extract_opendoc,
        }

        if suf in dispatch:
            text = dispatch[suf](path)
            if text:
                note = "Legacy PPT heuristic extraction" if suf == ".ppt" else None
                return text, True, note
            return "", True, f"Unable to extract {suf.lstrip('.').upper()} text"

        # Unknown extension – try heuristic before giving up
        text = self._extract_printable_runs(path)
        if text:
            return text, True, f"Unknown extension {suf!r}: heuristic extraction used"
        return "", False, f"Unhandled extension: {suf or '<none>'}"

    # ------------------------------------------------------------------
    # Per-format extractors
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_rtf(rtf: str) -> str:
        text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", rtf)
        text = re.sub(r"\\[a-zA-Z]+-?\d*\s?", " ", text)
        return text.replace("{", " ").replace("}", " ")

    @staticmethod
    def _extract_docx(path: Path) -> str:
        try:
            import docx  # type: ignore
            doc = docx.Document(str(path))
            parts: list[str] = []
            for para in doc.paragraphs:
                if para.text.strip():
                    style = para.style.name if para.style else ""
                    prefix = f"[{style}] " if style and style not in ("Normal", "Default Paragraph Style") else ""
                    parts.append(f"{prefix}{para.text}")
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            if parts:
                return "\n".join(parts)
        except Exception:
            pass
        # Fallback: raw XML
        try:
            with zipfile.ZipFile(path) as z:
                xml = z.read("word/document.xml")
            root = etree.fromstring(xml)
            return " ".join(root.itertext())
        except Exception:
            return ""

    @staticmethod
    def _extract_pdf(path: Path) -> str:
        # Try pypdf first
        try:
            from pypdf import PdfReader  # type: ignore
            reader = PdfReader(str(path))
            parts: list[str] = []
            for i, page in enumerate(reader.pages, 1):
                t = page.extract_text() or ""
                if t.strip():
                    parts.append(f"[Page {i}]\n{t}")
            if parts:
                return "\n".join(parts)
        except Exception:
            pass
        # Try pdfplumber
        try:
            import pdfplumber  # type: ignore
            with pdfplumber.open(str(path)) as pdf:
                parts = []
                for i, page in enumerate(pdf.pages, 1):
                    t = page.extract_text() or ""
                    if t.strip():
                        parts.append(f"[Page {i}]\n{t}")
                if parts:
                    return "\n".join(parts)
        except Exception:
            pass
        return TenderAgenticPipeline._extract_printable_runs(path)

    @staticmethod
    def _extract_doc(path: Path) -> str:
        for tool in ("antiword", "catdoc"):
            exe = shutil.which(tool)
            if exe:
                try:
                    r = subprocess.run(
                        [exe, str(path)],
                        capture_output=True, text=True,
                        encoding="utf-8", errors="ignore",
                        timeout=30,
                    )
                    if r.returncode == 0 and r.stdout.strip():
                        return r.stdout
                except Exception:
                    pass
        return TenderAgenticPipeline._extract_printable_runs(path)

    @staticmethod
    def _extract_xlsx(path: Path) -> str:
        try:
            from openpyxl import load_workbook  # type: ignore
            wb = load_workbook(str(path), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in wb.worksheets:
                parts.append(f"\n[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(min_row=1, max_row=500):
                    vals = [str(c.value).strip() for c in row if c.value is not None]
                    if vals:
                        parts.append(" | ".join(vals))
            return "\n".join(parts)
        except Exception:
            return ""

    @staticmethod
    def _extract_xls(path: Path) -> str:
        try:
            import xlrd  # type: ignore
            book = xlrd.open_workbook(str(path), on_demand=True)
            parts: list[str] = []
            for sheet in book.sheets():
                parts.append(f"\n[Sheet: {sheet.name}]")
                for r in range(min(sheet.nrows, 500)):
                    vals = [str(v).strip() for v in sheet.row_values(r) if str(v).strip()]
                    if vals:
                        parts.append(" | ".join(vals))
            return "\n".join(parts)
        except Exception:
            return ""

    @staticmethod
    def _extract_pptx(path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(str(path))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"\n[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts)
        except Exception:
            return ""

    @staticmethod
    def _extract_ppt_legacy(path: Path) -> str:
        return TenderAgenticPipeline._extract_printable_runs(path)

    @staticmethod
    def _extract_opendoc(path: Path) -> str:
        try:
            with zipfile.ZipFile(path) as z:
                xml = z.read("content.xml")
            root = etree.fromstring(xml)
            return " ".join(root.itertext())
        except Exception:
            return ""

    @staticmethod
    def _extract_printable_runs(path: Path) -> str:
        try:
            data = path.read_bytes()
        except OSError:
            return ""
        candidate = data.decode("latin-1", errors="ignore")
        runs = re.findall(r"[A-Za-z0-9][^\r\n]{20,}", candidate)
        return "\n".join(runs[:600])

    # ------------------------------------------------------------------
    # Text normalisation
    # ------------------------------------------------------------------

    @staticmethod
    def _normalize_text(text: str) -> str:
        if not text:
            return ""
        lower = text.lower()
        if "<html" in lower or "<table" in lower or "<body" in lower:
            try:
                tree = html.fromstring(text)
                text = tree.text_content()
            except Exception:
                pass
        lines = [l.strip() for l in text.splitlines()]
        out: list[str] = []
        blank_run = 0
        for l in lines:
            if l:
                blank_run = 0
                out.append(l)
            else:
                blank_run += 1
                if blank_run <= 2:
                    out.append("")
        return "\n".join(out).strip()

    # ------------------------------------------------------------------
    # Stage 1: triage
    # ------------------------------------------------------------------

    def _triage_relevant_files(
        self, tender_id: str, docs: list[FileDocument]
    ) -> list[str]:
        manifest = [
            {
                "file": d.rel_path,
                "ext": d.extension,
                "size_bytes": d.size_bytes,
                "chars": len(d.text),
                "sample": d.text[:500],
            }
            for d in docs
        ]

        messages = [
            {
                "role": "system",
                "content": (
                    "You are an expert procurement analyst. "
                    "Select the files most likely to contain:\n"
                    "- Award notices, contract award decisions\n"
                    "- Financial details: contract value, bid prices, BoQ, rates\n"
                    "- Technical specifications and scope of work\n"
                    "- Bidder/vendor details and evaluation scores\n"
                    "- Timeline, milestones, delivery schedules\n"
                    "- Terms and conditions, SLAs, penalties\n"
                    "- Amendments, corrigenda, addenda\n"
                    "- Eligibility and qualification criteria\n"
                    "Prefer larger, richer files. Include ALL spreadsheets (they often contain bid prices / BoQ). "
                    "Return ONLY valid JSON."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Tender ID: {tender_id}\n"
                    f"Select up to {self.max_relevant_files} files.\n"
                    'JSON: {"relevant_files":["path"],"reasoning":"..."}\n\n'
                    f"Manifest:\n{json.dumps(manifest, ensure_ascii=True)}"
                ),
            },
        ]

        try:
            raw = self.llm.chat(messages, temperature=0.0, max_tokens=1000)
            payload = _parse_json(raw)
            selected = [str(x) for x in payload.get("relevant_files", []) if isinstance(x, str)]
        except Exception:
            selected = []

        known = {d.rel_path for d in docs}
        selected = [s for s in selected if s in known]
        if not selected:
            selected = [d.rel_path for d in sorted(docs, key=lambda d: d.size_bytes, reverse=True)]
        return selected[: self.max_relevant_files]

    # ------------------------------------------------------------------
    # Stage 2: deep per-file evidence extraction
    # ------------------------------------------------------------------

    def _extract_deep_evidence(
        self,
        tender_id: str,
        docs: list[FileDocument],
        relevant_files: list[str],
    ) -> list[dict[str, Any]]:
        doc_map = {d.rel_path: d for d in docs}
        evidence: list[dict[str, Any]] = []

        for rel_path in relevant_files:
            doc = doc_map.get(rel_path)
            if not doc:
                continue
            self._log(f"  Extracting evidence: {rel_path}")

            schema = """{
  "file": "...",
  "document_type": "award_notice|bid_document|specification|contract|amendment|evaluation|schedule|correspondence|other",
  "status_hint": "awarded|not_awarded|cancelled|under_evaluation|unknown",

  "contracting_authority": {
    "name": "...", "department": "...", "address": "...",
    "contact_person": "...", "email": "...", "phone": "...",
    "evidence": "..."
  },

  "awarded_to": [
    {
      "name": "...", "registration_number": "...", "address": "...",
      "contact": "...", "email": "...", "consortium_members": ["..."],
      "evidence": "..."
    }
  ],

  "contract_value": {
    "total_amount": "...", "currency": "...",
    "breakdown": [{"item":"...","amount":"...","unit":"...","quantity":"..."}],
    "vat_rate": "...", "vat_amount": "...", "amount_ex_vat": "...",
    "evidence": "..."
  },

  "bid_submissions": [
    {
      "bidder": "...", "bid_amount": "...", "currency": "...",
      "technical_score": "...", "financial_score": "...", "total_score": "...",
      "rank": "...", "disqualified": false, "disqualification_reason": "...",
      "evidence": "..."
    }
  ],

  "timeline": [
    {"event": "...", "date": "...", "deadline": false, "evidence": "..."}
  ],

  "scope_of_work": {
    "description": "...",
    "deliverables": ["..."],
    "locations": ["..."],
    "duration": "...",
    "start_date": "...", "end_date": "...",
    "evidence": "..."
  },

  "technical_specifications": [
    {"spec": "...", "requirement": "...", "unit": "...", "value": "...", "evidence": "..."}
  ],

  "evaluation_criteria": [
    {"criterion": "...", "weight_pct": "...", "max_score": "...", "sub_criteria": ["..."], "evidence": "..."}
  ],

  "eligibility_criteria": [
    {"criterion": "...", "mandatory": true, "minimum_requirement": "...", "evidence": "..."}
  ],

  "sla_and_kpis": [
    {"metric": "...", "target": "...", "measurement_method": "...", "reporting_frequency": "...", "evidence": "..."}
  ],

  "penalties_and_ld": [
    {"trigger": "...", "rate": "...", "cap": "...", "grace_period": "...", "evidence": "..."}
  ],

  "payment_terms": {
    "payment_schedule": "...",
    "milestones": [{"milestone":"...","pct":"...","amount":"...","trigger":"..."}],
    "retention_pct": "...", "retention_release": "...",
    "advance_payment": "...", "advance_recovery": "...",
    "payment_period_days": "...",
    "evidence": "..."
  },

  "warranty_and_support": {
    "warranty_period": "...", "warranty_scope": "...",
    "support_terms": "...", "support_hours": "...",
    "evidence": "..."
  },

  "amendments_and_corrigenda": [
    {"amendment_no": "...", "date": "...", "description": "...", "impact": "...", "evidence": "..."}
  ],

  "legal_and_compliance": [
    {"requirement": "...", "law_reference": "...", "jurisdiction": "...", "evidence": "..."}
  ],

  "financial_security": {
    "bid_bond_amount": "...", "bid_bond_validity": "...",
    "performance_bond_pct": "...", "performance_bond_amount": "...",
    "retention_money_pct": "...", "advance_payment_guarantee": "...",
    "evidence": "..."
  },

  "subcontracting": {
    "allowed": null, "max_pct": "...", "approval_required": null,
    "requirements": "...", "evidence": "..."
  },

  "key_personnel": [
    {"role": "...", "name": "...", "qualification": "...", "min_experience_years": "...", "evidence": "..."}
  ],

  "insurance_requirements": [
    {"type": "...", "minimum_cover": "...", "evidence": "..."}
  ],

  "documents_required": ["..."],
  "documents_submitted": ["..."],

  "references_and_cross_links": ["..."],
  "extra_notes": ["..."]
}"""

            messages = [
                {
                    "role": "system",
                    "content": (
                        "You are a meticulous procurement document analyst. "
                        "Extract EVERY piece of useful information from the document. "
                        "Be exhaustive – capture all numbers, dates, names, amounts, requirements, and clauses. "
                        "Never summarise away detail. Quote exact phrases as evidence. "
                        "Return ONLY valid JSON matching the schema. "
                        "Omit top-level keys that are completely absent from this document."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"Tender ID: {tender_id}\nFile: {rel_path}\n\n"
                        f"Return JSON matching this schema:\n{schema}\n\n"
                        f"Document text:\n{doc.text}"
                    ),
                },
            ]

            try:
                raw = self.llm.chat(messages, temperature=0.0, max_tokens=3500)
                payload = _parse_json(raw)
                if isinstance(payload, dict):
                    evidence.append(payload)
            except Exception as exc:
                evidence.append({
                    "file": rel_path,
                    "status_hint": "unknown",
                    "error": f"Evidence extraction failed: {exc}",
                })

        return evidence

    # ------------------------------------------------------------------
    # Stage 3: deep synthesis
    # ------------------------------------------------------------------

    def _synthesize_detailed(
        self,
        tender_id: str,
        docs: list[FileDocument],
        evidence: list[dict[str, Any]],
    ) -> dict[str, Any]:
        manifest = [
            {"file": d.rel_path, "ext": d.extension, "chars": len(d.text)}
            for d in docs
        ]

        output_schema = """{
  "tender_id": "...",
  "status": "awarded|not_awarded|cancelled|under_evaluation|unknown",
  "confidence": 0.0,

  "summary": "Comprehensive 3-5 sentence plain English summary of the tender, outcome, winning vendor, contract value, and key terms.",

  "contracting_authority": {
    "name": "...", "department": "...", "address": "...",
    "contact_person": "...", "email": "...", "phone": "...",
    "evidence": ["..."]
  },

  "awarded_to": {
    "primary_contractor": "...",
    "registration_number": "...",
    "address": "...",
    "consortium_members": ["..."],
    "contact": "...", "email": "...",
    "evidence": ["..."]
  },

  "contract_value": {
    "total_amount": "...", "currency": "...",
    "amount_ex_vat": "...", "vat_rate": "...", "vat_amount": "...",
    "breakdown": [{"item":"...","amount":"...","unit":"...","quantity":"..."}],
    "evidence": ["..."]
  },

  "all_bidders": [
    {
      "name": "...", "bid_amount": "...", "currency": "...",
      "technical_score": "...", "financial_score": "...", "total_score": "...",
      "rank": "...", "outcome": "awarded|rejected|disqualified",
      "disqualification_reason": "...",
      "evidence": ["..."]
    }
  ],

  "scope_of_work": {
    "description": "...",
    "deliverables": ["..."],
    "locations": ["..."],
    "duration": "...",
    "start_date": "...", "end_date": "...",
    "evidence": ["..."]
  },

  "timeline": [
    {"event": "...", "date": "...", "deadline": false, "evidence": "..."}
  ],

  "technical_specifications": [
    {"spec": "...", "requirement": "...", "unit": "...", "value": "...", "evidence": ["..."]}
  ],

  "evaluation_criteria": [
    {"criterion": "...", "weight_pct": "...", "max_score": "...", "sub_criteria": ["..."], "evidence": ["..."]}
  ],

  "eligibility_criteria": [
    {"criterion": "...", "mandatory": true, "minimum_requirement": "...", "evidence": ["..."]}
  ],

  "sla_and_kpis": [
    {"metric": "...", "target": "...", "measurement_method": "...", "reporting_frequency": "...", "evidence": ["..."]}
  ],

  "penalties_and_ld": [
    {"trigger": "...", "rate": "...", "cap": "...", "grace_period": "...", "evidence": ["..."]}
  ],

  "payment_terms": {
    "payment_schedule": "...",
    "milestones": [{"milestone":"...","pct":"...","amount":"...","trigger":"..."}],
    "retention_pct": "...", "retention_release": "...",
    "advance_payment": "...", "advance_recovery": "...",
    "payment_period_days": "...",
    "evidence": ["..."]
  },

  "warranty_and_support": {
    "warranty_period": "...", "warranty_scope": "...",
    "support_terms": "...", "support_hours": "...",
    "evidence": ["..."]
  },

  "financial_security": {
    "bid_bond_amount": "...", "bid_bond_validity": "...",
    "performance_bond_pct": "...", "performance_bond_amount": "...",
    "retention_money_pct": "...", "advance_payment_guarantee": "...",
    "evidence": ["..."]
  },

  "amendments_and_corrigenda": [
    {"amendment_no":"...","date":"...","description":"...","impact":"...","evidence":["..."]}
  ],

  "legal_and_compliance": [
    {"requirement":"...","law_reference":"...","jurisdiction":"...","evidence":["..."]}
  ],

  "subcontracting": {
    "allowed": null, "max_pct": "...", "approval_required": null,
    "requirements": "...", "evidence": ["..."]
  },

  "insurance_requirements": [
    {"type": "...", "minimum_cover": "...", "evidence": ["..."]}
  ],

  "key_personnel": [
    {"role":"...","name":"...","qualification":"...","min_experience_years":"...","evidence":["..."]}
  ],

  "gaps": ["exhaustive list of every field or fact that is missing or unclear"],

  "source_conflicts": [
    {"field":"...","values":["..."],"resolution":"...","sources":["..."]}
  ],

  "sources_used": ["path1","path2"]
}"""

        messages = [
            {
                "role": "system",
                "content": (
                    "You are a senior procurement analyst producing an authoritative, "
                    "comprehensive tender dossier. "
                    "Merge ALL per-file evidence into one exhaustive structured report. "
                    "Resolve conflicts by preferring the most authoritative source "
                    "(award notice > evaluation report > contract > bid document). "
                    "Populate EVERY field you have data for. "
                    "Use null only when data is genuinely absent after searching all evidence. "
                    "The 'gaps' array must list every genuinely missing piece of information. "
                    "Return ONLY valid JSON."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Tender ID: {tender_id}\n\n"
                    f"Produce the final dossier JSON:\n{output_schema}\n\n"
                    f"File manifest:\n{json.dumps(manifest, ensure_ascii=True)}\n\n"
                    f"Evidence from all files:\n{json.dumps(evidence, ensure_ascii=True)}"
                ),
            },
        ]

        try:
            raw = self.llm.chat(messages, temperature=0.0, max_tokens=4000)
            payload = _parse_json(raw)
            if not isinstance(payload, dict):
                raise ValueError("Synthesis output is not a JSON object")
        except Exception as exc:
            payload = {
                "tender_id": tender_id,
                "status": "unknown",
                "confidence": 0.0,
                "summary": f"Synthesis failed: {exc}",
                "contracting_authority": None,
                "awarded_to": None,
                "contract_value": None,
                "all_bidders": [],
                "scope_of_work": None,
                "timeline": [],
                "technical_specifications": [],
                "evaluation_criteria": [],
                "eligibility_criteria": [],
                "sla_and_kpis": [],
                "penalties_and_ld": [],
                "payment_terms": None,
                "warranty_and_support": None,
                "financial_security": None,
                "amendments_and_corrigenda": [],
                "legal_and_compliance": [],
                "subcontracting": None,
                "insurance_requirements": [],
                "key_personnel": [],
                "gaps": [f"Synthesis failed: {exc}"],
                "source_conflicts": [],
                "sources_used": [],
            }

        payload.setdefault("tender_id", tender_id)
        payload["evidence_per_file"] = evidence
        return payload


# ---------------------------------------------------------------------------
# JSON parsing helper
# ---------------------------------------------------------------------------

def _parse_json(text: str) -> dict[str, Any]:
    candidate = text.strip()
    if candidate.startswith("```"):
        candidate = re.sub(r"^```[a-zA-Z0-9_\-]*\n?", "", candidate).strip()
        if candidate.endswith("```"):
            candidate = candidate[:-3].strip()
    try:
        return json.loads(candidate)
    except json.JSONDecodeError:
        pass
    match = re.search(r"\{[\s\S]*\}", text)
    if match:
        try:
            return json.loads(match.group(0))
        except json.JSONDecodeError:
            pass
    raise ValueError("No valid JSON object found in model response")


# ---------------------------------------------------------------------------
# Tender discovery
# ---------------------------------------------------------------------------

def discover_tenders(input_root: Path) -> list[tuple[str, Path]]:
    """
    Supports multiple layouts:
      <root>/<tender_id>/                        → use that dir directly
      <root>/<tender_id>/<tender_id>_extracted/
      <root>/<tender_id>/inspect/
    """
    tenders: list[tuple[str, Path]] = []
    if not input_root.exists():
        return tenders

    for tender_dir in sorted(input_root.iterdir(), key=lambda p: p.name):
        if not tender_dir.is_dir():
            continue
        if not re.fullmatch(r"\d{6}", tender_dir.name):
            continue

        extracted = sorted(tender_dir.glob("*_extracted"))
        if extracted:
            tenders.append((tender_dir.name, extracted[0]))
            continue

        inspect_dir = tender_dir / "inspect"
        if inspect_dir.is_dir():
            tenders.append((tender_dir.name, inspect_dir))
            continue

        # Use the tender dir itself – the pipeline will walk everything inside.
        tenders.append((tender_dir.name, tender_dir))

    return tenders


# ---------------------------------------------------------------------------
# I/O helpers
# ---------------------------------------------------------------------------

def write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


def append_jsonl(path: Path, payload: dict[str, Any], lock: threading.Lock) -> None:
    line = json.dumps(payload, ensure_ascii=False) + "\n"
    with lock:
        with path.open("a", encoding="utf-8") as fh:
            fh.write(line)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Deeply detailed tender agentic pipeline (Ollama default / vLLM).",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )

    p.add_argument("--input-root",  default="downloads/batch_runs")
    p.add_argument("--output-dir",  default="outputs/tender_structured")
    p.add_argument("--backend",     choices=["ollama", "vllm"], default="vllm")

    g_vl = p.add_argument_group("vLLM / OpenAI-compatible")
    g_vl.add_argument("--vllm-base-url", default=os.getenv("VLLM_BASE_URL", "http://10.4.25.56:8000/v1"))
    g_vl.add_argument("--vllm-model",    default=os.getenv("VLLM_MODEL",    "Qwen/Qwen3-30B-A3B-GPTQ-Int4"))
    g_vl.add_argument("--vllm-api-key",  default=os.getenv("VLLM_API_KEY",  "EMPTY"))
    g_vl.add_argument("--vllm-timeout",  type=float, default=180.0)

    p.add_argument("--workers",            type=int,   default=2)
    p.add_argument("--limit",              type=int,   default=0,    help="Process only first N tenders (0=all).")
    p.add_argument("--max-files",          type=int,   default=300,  help="Max files to read per tender.")
    p.add_argument("--max-file-chars",     type=int,   default=18000, help="Max characters to keep per file.")
    p.add_argument("--max-relevant-files", type=int,   default=20,   help="Max files passed to evidence extraction.")
    p.add_argument("--heartbeat-seconds",  type=float, default=30.0)
    p.add_argument("--no-unpack",          action="store_true", help="Skip archive unpacking.")
    p.add_argument("--verbose",            action="store_true")
    p.add_argument("--fail-on-unhandled",  action="store_true")

    return p


def _build_llm(args: argparse.Namespace) -> BaseLLMClient:
    if args.backend == "vllm":
        return VllmChatClient(
            base_url=args.vllm_base_url,
            model=args.vllm_model,
            api_key=args.vllm_api_key,
            timeout_seconds=args.vllm_timeout,
        )
    client = OllamaClient(
        base_url=args.ollama_base_url,
        model=args.ollama_model,
        timeout_seconds=args.ollama_timeout,
    )
    if not client.check_model():
        print(f"[WARN] Model '{args.ollama_model}' not found in Ollama. Run: ollama pull {args.ollama_model}")
    return client


def main() -> int:
    args = build_parser().parse_args()

    input_root = Path(args.input_root).resolve()
    output_dir = Path(args.output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    tenders = discover_tenders(input_root)
    if args.limit > 0:
        tenders = tenders[: args.limit]
    if not tenders:
        print(f"No tender folders found under: {input_root}")
        return 1

    llm = _build_llm(args)
    pipeline = TenderAgenticPipeline(
        llm,
        max_files=args.max_files,
        max_file_chars=args.max_file_chars,
        max_relevant_files=args.max_relevant_files,
        unpack_archives=not args.no_unpack,
        verbose=args.verbose,
    )

    all_results_path = output_dir / "all_tenders.jsonl"
    if all_results_path.exists():
        all_results_path.unlink()
    jsonl_lock = threading.Lock()

    print(f"Discovered tenders   : {len(tenders)}")
    print(f"Backend              : {args.backend}")
    if args.backend == "ollama":
        print(f"Ollama URL / model   : {args.ollama_base_url}  /  {args.ollama_model}")
    else:
        print(f"vLLM URL / model     : {args.vllm_base_url}  /  {args.vllm_model}")
    print(f"Workers              : {args.workers}")
    print(f"Max files per tender : {args.max_files}")
    print(f"Max chars per file   : {args.max_file_chars}")
    print(f"Max relevant files   : {args.max_relevant_files}")
    if args.limit > 0:
        print(f"Limit                : first {args.limit} tenders")

    total_issues = 0
    completed = 0
    started_at = time.monotonic()

    def run_one(item: tuple[str, Path]) -> tuple[str, dict[str, Any]]:
        tid, root = item
        result = pipeline.process_tender(tid, root)
        return tid, result

    with ThreadPoolExecutor(max_workers=max(1, args.workers)) as executor:
        futures = {executor.submit(run_one, item): item[0] for item in tenders}

        while futures:
            done_batch, _ = wait(
                set(futures.keys()),
                timeout=max(1.0, args.heartbeat_seconds),
                return_when=FIRST_COMPLETED,
            )
            if not done_batch:
                elapsed = time.monotonic() - started_at
                print(
                    f"[HEARTBEAT] completed={completed}/{len(tenders)} "
                    f"in-flight={len(futures)} elapsed={elapsed:.1f}s"
                )
                continue

            for future in done_batch:
                tid = futures.pop(future)
                try:
                    result_tid, result = future.result()
                except Exception as exc:
                    result_tid = tid
                    result = {
                        "tender_id": tid,
                        "status": "unknown",
                        "error": f"Pipeline error: {exc}",
                        "extraction_issues": [],
                    }

                per_path = output_dir / f"{result_tid}.json"
                write_json(per_path, result)
                append_jsonl(all_results_path, result, jsonl_lock)

                issues = result.get("extraction_issues", [])
                issue_count = len(issues) if isinstance(issues, list) else 0
                total_issues += issue_count
                completed += 1
                n_files = len(result.get("all_files_found", []))
                print(
                    f"[DONE] {result_tid} ({completed}/{len(tenders)})  "
                    f"files_read={n_files}  issues={issue_count}  -> {per_path}"
                )

    elapsed_total = time.monotonic() - started_at
    print(f"\nConsolidated JSONL   : {all_results_path}")
    print(f"Total elapsed        : {elapsed_total:.1f}s")
    if total_issues:
        print(f"Total extraction issues: {total_issues}")
        if args.fail_on_unhandled:
            return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())